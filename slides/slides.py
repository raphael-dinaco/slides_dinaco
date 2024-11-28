import sys
sys.path.append("..")
import pandas as pd
import datetime as dt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt, Cm
from Automação.modulos.Tratamento_bases import *
from Automação.modulos.Funcoes_slides import *


df_faturamento = tratamento_faturamento()
df_oportunidade = tratamento_oportunidade()
df_compras = tratamento_compras()
tabela_compras = tabela_compras()
df_visitas = tratamento_visitas()
df_meta = tratamento_meta()

#--------------------SLIDES DE INICIO E FIM--------------
def slide_inicio(prs_name, end_date, start_date, prs, representada):
    layout_name = 'dados_apresentacao'

    # Find the slide layout by name and add a new slide
    slide_layout = find_layout_by_name(prs, layout_name)
    slide = prs.slides.add_slide(slide_layout)

    # Set the title for the slide
    text_placeholder = slide.placeholders[0]  # Index 0 for the TITLE placeholder (available in "template_placeholder_mapping.txt)")
    text_placeholder.text = prs_name

    # Set the summary for the data
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = ('Data de Criação: ' + str(dt.datetime.now().strftime('%d/%m/%Y')) + '\n'
                            + 'intervalo: ' + str(start_date.strftime('%m/%Y')) + '-' + str(end_date.strftime('%m/%Y'))
                            + '\n' + 'Representada: ' + ', '.join(representada))
    
def slide_fim(prs):
    layout_name = 'fim_apresentacao'

    # Find the slide layout by name and add a new slide
    slide_layout = find_layout_by_name(prs, layout_name)
    slide = prs.slides.add_slide(slide_layout)

#--------------------SLIDES DE DIRETORIA-----------------
def crescimento_queda_top20(end_date, prs):
    df = df_faturamento[(df_faturamento['Mês/Ano'].dt.month <= end_date.month)&(df_faturamento['Mês/Ano'].dt.year>=(end_date.year-4))&(df_faturamento['Mês/Ano'].dt.year<=(end_date.year))][['GRUPOECONOMICO','ANO','VLRFAT_TOT_USD']].groupby(['GRUPOECONOMICO','ANO']).sum().reset_index()
    df['VLRFAT_TOT_USD'] = round(df['VLRFAT_TOT_USD'], 0)
    df = df.pivot(index='GRUPOECONOMICO', columns='ANO', values='VLRFAT_TOT_USD').fillna(0).reset_index()
    df['Dif. 24-23'] = df[end_date.year] - df[end_date.year-1]
    df = df.merge(df_visitas[(df_visitas['AD_DHVISITA'].dt.year == end_date.year) & 
            (df_visitas['AD_DHVISITA'].dt.month <= end_date.month)
            ][['AD_GRUPOECONOMICO', 'DESCRHIST']].groupby(['AD_GRUPOECONOMICO', 'DESCRHIST']).size().reset_index().pivot(
                columns='DESCRHIST', index='AD_GRUPOECONOMICO', values=0
                ).fillna(0).reset_index(), left_on='GRUPOECONOMICO', right_on='AD_GRUPOECONOMICO', how='left').fillna(0)
    df = df.drop(columns=['AD_GRUPOECONOMICO'])
    df = df.rename(columns={'GRUPOECONOMICO': 'CLIENTE'})
    df_top = df.sort_values(by= 'Dif. 24-23', ascending=False).head(20)
    df_down = df.sort_values(by= 'Dif. 24-23', ascending=True).head(20)
    df_top.iloc[:, -8:-3] = df_top.iloc[:, -8:-3].applymap(lambda x: '${:,.0f}'.format(x)).replace(',', '.', regex=True)
    df_top.iloc[:, -3:] = df_top.iloc[:, -3:].applymap(lambda x: '{:,.0f}'.format(x)).replace(',', '.', regex=True)
    df_down.iloc[:, -8:-3] = df_down.iloc[:, -8:-3].applymap(lambda x: '${:,.0f}'.format(x)).replace(',', '.', regex=True)
    df_down.iloc[:, -3:] = df_down.iloc[:, -3:].applymap(lambda x: '{:,.0f}'.format(x)).replace(',', '.', regex=True)
    

    #-------------------SLIDE-------------------#
    # Get the layout of the slide
    layout_name = 'tabela_sem_coment'

    # Find the slide layout by name and add a new slide
    slide_layout = find_layout_by_name(prs, layout_name)
    slide = prs.slides.add_slide(slide_layout)

    # Set the title for the slide
    text_placeholder = slide.placeholders[0]  # Index 0 for the TITLE placeholder (available in "template_placeholder_mapping.txt)")
    text_placeholder.text = 'TOP 20 - Maior crescimento YTD'

    individual_width = [12.18,2.9,2.9,2.9,2.9,2.9,2.92,1.65,1.8]
    nova_tabela(df_top, slide, individual_width,0.42,1.99,0,11.5)

    #-------------------------------SLIDE DE CLIENTES COM MAIOR QUEDA DE FATURAMENTO---------------------------------#

    # Get the layout of the slide
    layout_name = 'tabela_sem_coment'

    # Find the slide layout by name and add a new slide
    slide_layout = find_layout_by_name(prs, layout_name)
    slide = prs.slides.add_slide(slide_layout)

    # Set the title for the slide
    text_placeholder = slide.placeholders[0]  # Index 0 for the TITLE placeholder (available in "template_placeholder_mapping.txt)")
    text_placeholder.text = 'TOP 20 - Maior queda YTD'

    individual_width = [12.18,2.9,2.9,2.9,2.9,2.9,2.92,1.65,1.8]
    nova_tabela(df_down, slide, individual_width, 0.42, 1.99,0,11.5)

def sem_faturamento (end_date, prs, lim_table):
    #Racional - Clientes que não tiveram faturamento no período determinado, mas tiveram faturamento no mesmo período do ano anterior
    bus = [['LIFE & PERSONAL CARE', 'FOOD & NUTRITION'],['INDUSTRIAL','AGRO']]
    abreviacoes = ['LPC/FN', 'IND/AGRO']
    for bu, abreviacao in zip(bus, abreviacoes):
        df = df_faturamento[
            (df_faturamento['Mês/Ano'].dt.year>=(end_date.year-5))&
            (df_faturamento['Mês/Ano'].dt.year<=(end_date.year))&
            (df_faturamento['UNIDADENEG'].isin(bu))
            ][['GRUPOECONOMICO','ANO','VLRFAT_TOT_USD']].groupby(['GRUPOECONOMICO','ANO']).sum().reset_index()
        df['VLRFAT_TOT_USD'] = round(df['VLRFAT_TOT_USD'], 0)
        df = df.pivot(index='GRUPOECONOMICO', columns='ANO', values='VLRFAT_TOT_USD').fillna(0).reset_index()
        df = df[(df[end_date.year] <= 0) & (df[end_date.year-1] > 0)]
        df = df.merge(df_visitas[(df_visitas['AD_DHVISITA'].dt.year == end_date.year) & 
                (df_visitas['AD_DHVISITA'].dt.month <= end_date.month)
                ][['AD_GRUPOECONOMICO', 'DESCRHIST']].groupby(['AD_GRUPOECONOMICO', 'DESCRHIST']).size().reset_index().pivot(
                    columns='DESCRHIST', index='AD_GRUPOECONOMICO', values=0
                    ).fillna(0).reset_index(), left_on='GRUPOECONOMICO', right_on='AD_GRUPOECONOMICO', how='left').fillna(0)
        df = df.drop(columns=['AD_GRUPOECONOMICO', end_date.year])
        df = df.sort_values(by= end_date.year-1, ascending=False).head(lim_table)
        df.iloc[:, -7:-2] = df.iloc[:, -7:-2].applymap(lambda x: '${:,.0f}'.format(x)).replace(',', '.', regex=True)
        df.iloc[:, -2:] = df.iloc[:, -2:].applymap(lambda x: '{:,.0f}'.format(x)).replace(',', '.', regex=True)
        df = df.rename(columns={'GRUPOECONOMICO': 'CLIENTE'})
        #-------------------SLIDE-------------------#
        # Get the layout of the slide
        layout_name = 'tabela_sem_coment'

        # Find the slide layout by name and add a new slide
        slide_layout = find_layout_by_name(prs, layout_name)
        slide = prs.slides.add_slide(slide_layout)

        # Set the title for the slide
        text_placeholder = slide.placeholders[0]  # Index 0 for the TITLE placeholder (available in "template_placeholder_mapping.txt)")
        text_placeholder.text = 'Clientes sem faturamento YTD - ' + abreviacao

        individual_width = [13.33,2.85,2.85,2.85,2.85,2.85,2.85,2.29,2.2]
        nova_tabela(df, slide, individual_width, 0.7, 1.99)

def oport_em_aberto_gerencial (end_date, prs, group, lim_table): #Previstas para o ano vigente
    bus = [['LIFE & PERSONAL CARE', 'FOOD & NUTRITION'],['INDUSTRIAL','AGRO']]
    abreviacoes = ['LPC/FN', 'IND/AGRO']
    for bu, abreviacao in zip (bus, abreviacoes):
        #----------------CRIAÇÃO DO DATAFRAME DE REFERÊNCIA----------------#
        #Oportunidades com previsão para o ano de referência
        df = df_oportunidade[
            (df_oportunidade['DTESTFECHAMENTO'].dt.year == end_date.year) &
            (df_oportunidade['STATUS RESUMIDO'] == 'Em Andamento') &
            (df_oportunidade['UNIDADE_NEGOCIO_PAR'].isin(bu))
            ][['REPRESENTADA' ,'GRUPO_ECONOMICO','GRUPO_PRODUTO','NOME_PROJETO','DTESTFECHAMENTO','VLRTOT','QTDNEG']].sort_values(by='VLRTOT', ascending=False)

        total_projetos = str(df.shape[0])
        if group == 1:
            df = df.groupby(['GRUPO_ECONOMICO']).agg(
                {'NOME_PROJETO': lambda x : '/ '.join(set(str(i) for i in x)),
                 'GRUPO_PRODUTO': lambda x: ', '.join(set(str(i) for i in x)),
                 'REPRESENTADA': lambda x: ', '.join(set(str(i) for i in x)),
                 'QTDNEG': 'sum','VLRTOT':'sum', 'DTESTFECHAMENTO': 'max'
                 }).reset_index().sort_values(by='VLRTOT', ascending=False)

        df = df[['REPRESENTADA' ,'GRUPO_ECONOMICO','GRUPO_PRODUTO','NOME_PROJETO','DTESTFECHAMENTO','VLRTOT','QTDNEG']]
        df.rename(columns={'REPRESENTADA':'Representada','GRUPO_ECONOMICO': 'Cliente', 'NOME_PROJETO': 'Projeto', 'GRUPO_PRODUTO': 'Produto', 'QTDNEG': 'KG\n12 meses', 'VLRTOT': 'USD\n12 meses', 'DTESTFECHAMENTO': 'Previsão de Fechamento'}, inplace=True)
        df['Previsão de Fechamento'] = df['Previsão de Fechamento'].dt.strftime('%d/%m/%Y')
        peso_total = df['KG\n12 meses'].sum()
        dolar_total = df['USD\n12 meses'].sum()
        if df.shape[0] > lim_table:
            df = df.head(lim_table)
        df['USD\n12 meses'] = df['USD\n12 meses'].apply(lambda x: '${:,.0f}'.format(x)).astype(str).str.replace(',', '.')
        df['KG\n12 meses'] = df['KG\n12 meses'].apply(lambda x: '{:,.0f}'.format(x)).astype(str).str.replace(',', '.')
        peso_total = '{:,.0f}'.format(peso_total).replace(',', '.')
        dolar_total = '${:,.0f}'.format(dolar_total).replace(',', '.')

        #-------------------SLIDE-------------------#
        # Get the layout of the slide
        if group == 1:
            layout_name = 'tabela_sem_coment'
        else:
            layout_name = 'tabela_simples'

        # Find the slide layout by name and add a new slide
        slide_layout = find_layout_by_name(prs, layout_name)
        slide = prs.slides.add_slide(slide_layout)

        # Set the title for the slide
        text_placeholder = slide.placeholders[0]  # Index 0 for the TITLE placeholder (available in "template_placeholder_mapping.txt)")
        text_placeholder.text = 'Oportunidades previstas para ' + str(end_date.year) + ' - ' + abreviacao

        # Set the summary for the data
        if group == 1:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = total_projetos + 'projetos' + '\n' + str(peso_total) + ' Kg' + '\n' + str(dolar_total) + ' USD'
        else:
            text_placeholder = slide.placeholders[14]
            text_placeholder.text = total_projetos + 'projetos/ ' + str(peso_total) + ' Kg/ ' + str(dolar_total) + ' USD'

        individual_width = [3.67,3.87,8.32,9.28,2.75,2.3,2.31]
        nova_tabela(df, slide, individual_width,0.68,2.03,1,11.5)

def oport_abertas_periodo (end_date, start_date, prs, group, mes, lim_table):
    bus = [['LIFE & PERSONAL CARE', 'FOOD & NUTRITION'],['INDUSTRIAL','AGRO']]
    abreviacoes = ['LPC/FN', 'IND/AGRO']
    for bu, abreviacao in zip(bus, abreviacoes):
        #----------------CRIAÇÃO DO DATAFRAME DE REFERÊNCIA----------------#
        #Oportunidades com previsão para o ano de referência
        if mes == 1:
            df = df_oportunidade[
                (df_oportunidade['DTNEG'].dt.year == end_date.year) &
                (df_oportunidade['DTNEG'].dt.month == end_date.month) &
                (df_oportunidade['UNIDADE_NEGOCIO_PAR'].isin(bu))
                ][['REPRESENTADA' ,'GRUPO_ECONOMICO','GRUPO_PRODUTO','NOME_PROJETO','DTESTFECHAMENTO','VLRTOT','QTDNEG']]
        else:
            df = df_oportunidade[
                (df_oportunidade['DTNEG'].dt.year == end_date.year) &
                (df_oportunidade['DTNEG'].dt.month <= end_date.month) &
                (df_oportunidade['DTNEG'].dt.month >= start_date.month) &
                (df_oportunidade['UNIDADE_NEGOCIO_PAR'].isin(bu))
                ][['REPRESENTADA' ,'GRUPO_ECONOMICO','GRUPO_PRODUTO','NOME_PROJETO','DTESTFECHAMENTO','VLRTOT','QTDNEG']]

        total_projetos = str(df.shape[0])
        if group == 1:
            df = df.groupby(['GRUPO_ECONOMICO']).agg(
                {'NOME_PROJETO': lambda x : '/ '.join(set(str(i) for i in x)),
                 'GRUPO_PRODUTO': lambda x: ', '.join(set(str(i) for i in x)),
                 'REPRESENTADA': lambda x: ', '.join(set(str(i) for i in x)),
                 'QTDNEG': 'sum','VLRTOT':'sum', 'DTESTFECHAMENTO': 'max'
                 }).reset_index()

        df = df[['REPRESENTADA' ,'GRUPO_ECONOMICO','GRUPO_PRODUTO','NOME_PROJETO','DTESTFECHAMENTO','VLRTOT','QTDNEG']].sort_values(by='VLRTOT', ascending=False)
        df.rename(columns={'REPRESENTADA':'Representada','GRUPO_ECONOMICO': 'Cliente', 'NOME_PROJETO': 'Projeto', 'GRUPO_PRODUTO': 'Produto', 'QTDNEG': 'KG\n12 meses', 'VLRTOT': 'USD\n12 meses', 'DTESTFECHAMENTO': 'Previsão de Fechamento'}, inplace=True)
        df['Previsão de Fechamento'] = df['Previsão de Fechamento'].dt.strftime('%d/%m/%Y')
        peso_total = df['KG\n12 meses'].sum()
        dolar_total = df['USD\n12 meses'].sum()
        if df.shape[0] > lim_table:
            df = df.head(lim_table)
        df['USD\n12 meses'] = df['USD\n12 meses'].apply(lambda x: '${:,.0f}'.format(x)).astype(str).str.replace(',', '.')
        df['KG\n12 meses'] = df['KG\n12 meses'].apply(lambda x: '{:,.0f}'.format(x)).astype(str).str.replace(',', '.')
        peso_total = '{:,.0f}'.format(peso_total).replace(',', '.')
        dolar_total = '${:,.0f}'.format(dolar_total).replace(',', '.')

        #-------------------SLIDE-------------------#
        # Get the layout of the slide
        if group == 1:
            layout_name = 'tabela_sem_coment'
        else:
            layout_name = 'tabela_simples'

        # Find the slide layout by name and add a new slide
        slide_layout = find_layout_by_name(prs, layout_name)
        slide = prs.slides.add_slide(slide_layout)

        # Set the title for the slide
        text_placeholder = slide.placeholders[0]  # Index 0 for the TITLE placeholder (available in "template_placeholder_mapping.txt)")
        text_placeholder.text = 'Oportunidades abertas - ' + abreviacao

        # Set the summary for the data
        if group == 1:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = total_projetos + 'projetos' + '\n' + str(peso_total) + ' Kg' + '\n' + str(dolar_total)
        else:
            text_placeholder = slide.placeholders[14]
            text_placeholder.text = total_projetos + ' projetos/ ' + str(peso_total) + ' Kg/ ' + str(dolar_total)

        individual_width = [3.67,3.87,8.32,9.28,2.75,2.3,2.31]
        nova_tabela(df, slide, individual_width,0.68,1.99,1,11.5)

def oport_convertidas_gerencial (end_date, start_date, prs, group, mes, lim_table):
    bus = [['LIFE & PERSONAL CARE', 'FOOD & NUTRITION'],['INDUSTRIAL','AGRO']]
    abreviacoes = ['LPC/FN', 'IND/AGRO']

    for bu, abreviacao in zip(bus, abreviacoes):
        #----------------CRIAÇÃO DO DATAFRAME DE REFERÊNCIA----------------#
        if mes == 1:
            df = df_oportunidade[
                (df_oportunidade['DTFECHAMENTO'].dt.month == end_date.month) &
                (df_oportunidade['DTFECHAMENTO'].dt.year == end_date.year) &
                (df_oportunidade['STATUS RESUMIDO'].isin(['Faturado'])) &
                (df_oportunidade['UNIDADE_NEGOCIO_PAR'].isin(bu))
                ][['REPRESENTADA' ,'GRUPO_ECONOMICO','GRUPO_PRODUTO','NOME_PROJETO','VLRTOT','QTDNEG']]
        else:
            df = df_oportunidade[
                (df_oportunidade['DTFECHAMENTO'].dt.month <= end_date.month) &
                (df_oportunidade['DTFECHAMENTO'].dt.month >= start_date.month) &
                (df_oportunidade['DTFECHAMENTO'].dt.year == end_date.year) &
                (df_oportunidade['STATUS RESUMIDO'].isin(['Faturado'])) &
                (df_oportunidade['UNIDADE_NEGOCIO_PAR'].isin(bu))
                ][['REPRESENTADA' ,'GRUPO_ECONOMICO','GRUPO_PRODUTO','NOME_PROJETO','VLRTOT','QTDNEG']]


        total_projetos = str(df.shape[0])

        if group == 1:

            df = df.groupby(['GRUPO_ECONOMICO']).agg(
                {'NOME_PROJETO': lambda x : '/ '.join(set(str(i) for i in x)),
                 'GRUPO_PRODUTO': lambda x: ', '.join(set(str(i) for i in x)),
                 'REPRESENTADA': lambda x: ', '.join(set(str(i) for i in x)),
                 'QTDNEG': 'sum','VLRTOT':'sum'}).reset_index()
        
        df = df[['REPRESENTADA' ,'GRUPO_ECONOMICO','GRUPO_PRODUTO','NOME_PROJETO','VLRTOT','QTDNEG']].sort_values(by='VLRTOT', ascending=False)
        df.rename(columns={'REPRESENTADA':'Representada','GRUPO_ECONOMICO': 'Cliente', 'NOME_PROJETO': 'Projeto', 'GRUPO_PRODUTO': 'Produto', 'QTDNEG': 'KG\n12 meses', 'VLRTOT': 'USD\n12 meses', 'UNIDADE_NEGOCIO_PAR': 'Unidade de Negócio'}, inplace=True)
        peso_total = df['KG\n12 meses'].sum()
        dolar_total = df['USD\n12 meses'].sum()
        df['USD\n12 meses'] = df['USD\n12 meses'].apply(lambda x: '${:,.0f}'.format(x)).astype(str).str.replace(',', '.')
        df['KG\n12 meses'] = df['KG\n12 meses'].apply(lambda x: '{:,.0f}'.format(x)).astype(str).str.replace(',', '.')
        peso_total = '{:,.0f}'.format(peso_total).replace(',', '.')
        dolar_total = '${:,.0f}'.format(dolar_total).replace(',', '.')
        if df.shape[0] > lim_table:
            df = df.head(lim_table)

        #-------------------SLIDE-------------------#
        # Get the layout of the slide
        if group == 1:
            layout_name = 'tabela_sem_coment'
        else:
            layout_name = 'tabela_simples'

        # Find the slide layout by name and add a new slide
        slide_layout = find_layout_by_name(prs, layout_name)
        slide = prs.slides.add_slide(slide_layout)

        # Set the title for the slide
        text_placeholder = slide.placeholders[0]  # Index 0 for the TITLE placeholder (available in "template_placeholder_mapping.txt)")
        text_placeholder.text = 'Oportunidades faturadas - ' + abreviacao

        if group == 1:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = total_projetos + 'projetos' + '\n' + str(peso_total) + ' Kg' + '\n' + 'Valor total: ' + '\n' + str(dolar_total)
        else:
            text_placeholder = slide.placeholders[14]
            text_placeholder.text = total_projetos + ' projetos/ ' + str(peso_total) + ' Kg/ ' + 'Valor total: ' + str(dolar_total)

        individual_width = [4.75,6.25,8.30,6.50,3.12,3.28]
        nova_tabela(df, slide, individual_width,0.83,1.99,1,11.5)

def oport_sem_amostra_gerencial (end_date, prs):
    bus = [['LIFE & PERSONAL CARE', 'FOOD & NUTRITION'],['INDUSTRIAL','AGRO']]
    for bu in bus:
        #----------------CRIAÇÃO DO DATAFRAME DE REFERÊNCIA----------------#
        #Oportunidades com previsão para o ano de referência
        df_sujo = pd.read_csv('SQL/Outputs/Oportunidades.csv')
        df_sujo['DTNEG'] = pd.to_datetime(df_sujo['DTNEG'])
        df_sujo['DT_FATUR_AMOSTRA'] = pd.to_datetime(df_sujo['DT_FATUR_AMOSTRA'])
        df_sujo = df_sujo[(df_sujo['UNIDADE_NEGOCIO_PAR'].isin(bu)) &
            (df_sujo['STATUS'].isin(['Início de Projeto','Teste não iniciado','Aprovado','Teste de estabilidade','Teste em laboratório/formulação','Negociação','Teste piloto','Teste no consumidor'])) &
            (df_sujo['DTNEG'] <= (end_date + pd.DateOffset(months=1)))                           
        ][['REPRESENTADA' ,'GRUPO_ECONOMICO','GRUPO_PRODUTO','NOME_PROJETO','DTNEG','VLRTOT','QTDNEG','STATUS_AMOSTRA','DTESTFECHAMENTO','DT_FATUR_AMOSTRA']].sort_values(by='VLRTOT', ascending=False)
        total_amostras = str(df_sujo.shape[0])
        df_sujo_status = df_sujo.groupby('STATUS_AMOSTRA').agg({'STATUS_AMOSTRA':'count'}).rename(columns={'STATUS_AMOSTRA': 'N° AMOSTRAS'}).reset_index()

        df = df_oportunidade[
            (df_oportunidade['STATUS RESUMIDO'] == 'Em Andamento') &
            (df_oportunidade['UNIDADE_NEGOCIO_PAR'].isin(bu)) &
            (df_oportunidade['DTNEG'] <= (end_date + pd.DateOffset(months=1)))
            ][['REPRESENTADA' ,'GRUPO_ECONOMICO','GRUPO_PRODUTO','NOME_PROJETO','DTNEG','VLRTOT','QTDNEG','STATUS_AMOSTRA','DTESTFECHAMENTO','DT_FATUR_AMOSTRA']].sort_values(by='VLRTOT', ascending=False)
        total_projetos = str(df.shape[0])
        peso_total = df['QTDNEG'].sum()
        dolar_total = df['VLRTOT'].sum()
        peso_total = '{:,.0f}'.format(peso_total).replace(',', '.')
        dolar_total = '${:,.0f}'.format(dolar_total).replace(',', '.')
        df['Tempo de Espera'] = (dt.datetime.now() - df['DTNEG']).dt.days
        df_sujo['Tempo de Espera'] = (dt.datetime.now() - df_sujo['DTNEG']).dt.days
        df['Intervalo de Espera'] = pd.cut(df['Tempo de Espera'], bins=[0, 30, 60, 90, 120, float('inf')], labels=['0-30', '30-60', '60-90', '90-120', '120+'])
        df_sujo['Intervalo de Espera'] = pd.cut(df_sujo['Tempo de Espera'], bins=[0, 30, 60, 90, 120, float('inf')], labels=['0-30', '30-60', '60-90', '90-120', '120+'])
        df['Tempo de Envio'] = (dt.datetime.now() - df['DT_FATUR_AMOSTRA']).dt.days
        df_sujo['Tempo de Envio'] = (dt.datetime.now() - df_sujo['DT_FATUR_AMOSTRA']).dt.days
        df['Intervalo de Envio'] = pd.cut(df['Tempo de Envio'], bins=[0, 60, 120, 180, 240, 360, float('inf')], labels=['0-60', '60-120', '120-180', '180-240', '240-360', '360+'])
        df_sujo['Intervalo de Envio'] = pd.cut(df_sujo['Tempo de Envio'], bins=[0, 60, 120, 180, 240, 360, float('inf')], labels=['0-60', '60-120', '120-180', '180-240', '240-360', '360+'])
        df_status = df.groupby('STATUS_AMOSTRA').agg({'STATUS_AMOSTRA':'count', 'VLRTOT':'sum', 'QTDNEG':'sum'}).rename(columns={'STATUS_AMOSTRA': 'N° PROJETOS', 'VLRTOT': 'VALOR', 'QTDNEG': 'PESO'}).reset_index()
        df_status['VALOR'] = df_status['VALOR'].apply(lambda x: '${:,.0f}'.format(x).replace(',', '.'))
        df_status['PESO'] = df_status['PESO'].apply(lambda x: '{:,.0f}'.format(x).replace(',', '.'))
        df_status = df_status.merge(df_sujo_status, on='STATUS_AMOSTRA', how='left')
        df_status = df_status[['STATUS_AMOSTRA', 'N° PROJETOS','N° AMOSTRAS' ,'VALOR', 'PESO']]


        #------------------- DATAFRAME SEM AMOSTRA -------------------#
        df_sem_amostra = df[df['STATUS_AMOSTRA'] == '02 - Aguardando análise Logística/disp. Estoque']
        df_sujo_sem_amostra = df_sujo[df_sujo['STATUS_AMOSTRA'] == '02 - Aguardando análise Logística/disp. Estoque']
        total_sem_amostra = str(df_sem_amostra.shape[0])
        total_amostra_sem_amostra = str(df_sujo_sem_amostra.shape[0])
        peso_sem_amostra = df_sem_amostra['QTDNEG'].sum()
        dolar_sem_amostra = df_sem_amostra['VLRTOT'].sum()
        peso_sem_amostra = '{:,.0f}'.format(peso_sem_amostra).replace(',', '.')
        dolar_sem_amostra = '${:,.0f}'.format(dolar_sem_amostra).replace(',', '.')
        df_sem_amostra_intervalos = df_sem_amostra.groupby('Intervalo de Espera').agg({'Intervalo de Espera':'count', 'VLRTOT':'sum', 'QTDNEG':'sum'}).rename(columns={'Intervalo de Espera': 'N° PROJETOS', 'VLRTOT': 'VALOR', 'QTDNEG': 'PESO'}).reset_index()
        df_sujo_sem_amostra_intervalos = df_sujo_sem_amostra.groupby('Intervalo de Espera').agg({'Intervalo de Espera':'count'}).rename(columns={'Intervalo de Espera': 'N° AMOSTRAS'}).reset_index()
        df_sem_amostra_intervalos = df_sem_amostra_intervalos.merge(df_sujo_sem_amostra_intervalos, on='Intervalo de Espera', how='left')
        df_sem_amostra_intervalos['VALOR'] = df_sem_amostra_intervalos['VALOR'].apply(lambda x: '${:,.0f}'.format(x).replace(',', '.'))
        df_sem_amostra_intervalos['PESO'] = df_sem_amostra_intervalos['PESO'].apply(lambda x: '{:,.0f}'.format(x).replace(',', '.'))
        df_sem_amostra_intervalos = df_sem_amostra_intervalos[['Intervalo de Espera', 'N° PROJETOS', 'N° AMOSTRAS', 'VALOR', 'PESO']]

        #------------------- DATAFRAME COM AMOSTRA -------------------#
        df_com_amostra = df[df['STATUS_AMOSTRA'].isin(['05 - Amostra Enviada','06 - Amostra c/ código de rastreio'])]
        df_sujo_com_amostra = df_sujo[df_sujo['STATUS_AMOSTRA'].isin(['05 - Amostra Enviada','06 - Amostra c/ código de rastreio'])]
        total_com_amostra = str(df_com_amostra.shape[0])
        total_amostra_com_amostra = str(df_sujo_com_amostra.shape[0])
        peso_com_amostra = df_com_amostra['QTDNEG'].sum()
        dolar_com_amostra = df_com_amostra['VLRTOT'].sum()
        peso_com_amostra = '{:,.0f}'.format(peso_com_amostra).replace(',', '.')
        dolar_com_amostra = '${:,.0f}'.format(dolar_com_amostra).replace(',', '.')
        df_com_amostra_intervalos = df_com_amostra.groupby('Intervalo de Envio').agg({'Intervalo de Envio':'count', 'VLRTOT':'sum', 'QTDNEG':'sum'}).rename(columns={'Intervalo de Envio': 'N° PROJETOS', 'VLRTOT': 'VALOR', 'QTDNEG': 'PESO'}).reset_index()
        df_sujo_com_amostra_intervalos = df_sujo_com_amostra.groupby('Intervalo de Envio').agg({'Intervalo de Envio':'count'}).rename(columns={'Intervalo de Envio': 'N° AMOSTRAS'}).reset_index()
        df_com_amostra_intervalos = df_com_amostra_intervalos.merge(df_sujo_com_amostra_intervalos, on='Intervalo de Envio', how='left')
        df_com_amostra_intervalos['VALOR'] = df_com_amostra_intervalos['VALOR'].apply(lambda x: '${:,.0f}'.format(x).replace(',', '.'))
        df_com_amostra_intervalos['PESO'] = df_com_amostra_intervalos['PESO'].apply(lambda x: '{:,.0f}'.format(x).replace(',', '.'))
        df_com_amostra_intervalos = df_com_amostra_intervalos[['Intervalo de Envio', 'N° PROJETOS', 'N° AMOSTRAS', 'VALOR', 'PESO']]

        #-------------------SLIDE 1 - Oportunidade por Status-------------------#
        # Get the layout of the slide
        layout_name = 'tabela_simples'

        # Find the slide layout by name and add a new slide
        slide_layout = find_layout_by_name(prs, layout_name)
        slide = prs.slides.add_slide(slide_layout)

        # Set the title for the slide
        text_placeholder = slide.placeholders[0]  # Index 0 for the TITLE placeholder (available in "template_placeholder_mapping.txt)")
        text_placeholder.text = 'Oportunidades por Status de Amostra:' + str(end_date.month) + '/' + str(end_date.year) + ' - ' + bu[0]# + ' e ' + bu[1]

        # Set the summary for the data
        text_placeholder = slide.placeholders[14]  # Index 0 for the TITLE placeholder (available in "template_placeholder_mapping.txt)")
        text_placeholder.text = total_projetos + ' projetos' + '\n' + total_amostras + ' amostras' +'\n' + str(peso_total) + ' Kg' + '\n' + str(dolar_total)

        nova_tabela(df_status, slide, 23.84,1.5,3.2)

        #-------------------SLIDE 2 - Intervalos de Espera - Oportunidades em Andamento-------------------#
        # Get the layout of the slide
        layout_name = 'tabela_simples'

        # Find the slide layout by name and add a new slide
        slide_layout = find_layout_by_name(prs, layout_name)
        slide = prs.slides.add_slide(slide_layout)

        # Set the title for the slide
        text_placeholder = slide.placeholders[0]  # Index 0 for the TITLE placeholder (available in "template_placeholder_mapping.txt)")
        text_placeholder.text = 'Amostras Não Enviadas - Intervalos de Espera:' + str(end_date.month) + '/' + str(end_date.year) + ' - ' + bu[0] # + ' e ' + bu[1]

        # Set the summary for the data
        text_placeholder = slide.placeholders[14]  # Index 0 for the TITLE placeholder (available in "template_placeholder_mapping.txt)")
        text_placeholder.text = total_sem_amostra + ' projetos' +'\n' + total_amostra_sem_amostra + ' amostras' + '\n' + str(peso_sem_amostra) + ' Kg' + '\n' + str(dolar_sem_amostra) + '\n' + 'Intervalo médio de espera: ' + str(round(df_sujo_sem_amostra['Tempo de Espera'].mean(),0)) + ' dias'

        nova_tabela(df_sem_amostra_intervalos, slide, 23.84,1.5,3.2)

        #-------------------SLIDE 3 - Intervalo de Envio - Oportunidades em Andamento-------------------#
        # Get the layout of the slide
        layout_name = 'tabela_simples'

        # Find the slide layout by name and add a new slide
        slide_layout = find_layout_by_name(prs, layout_name)
        slide = prs.slides.add_slide(slide_layout)

        # Set the title for the slide
        text_placeholder = slide.placeholders[0]  # Index 0 for the TITLE placeholder (available in "template_placeholder_mapping.txt)")
        text_placeholder.text = 'Amostras Enviadas - Intervalos de Envio:' + str(end_date.month) + '/' + str(end_date.year) + ' - ' + bu[0]# + ' e ' + bu[1]

        # Set the summary for the data
        text_placeholder = slide.placeholders[14]  # Index 0 for the TITLE placeholder (available in "template_placeholder_mapping.txt)")
        text_placeholder.text = total_com_amostra + ' projetos' + '\n' + total_amostra_com_amostra + ' Amostras' + '\n' + str(peso_com_amostra) + ' Kg' + '\n' + str(dolar_com_amostra)

        nova_tabela(df_com_amostra_intervalos, slide, 23.84,1.5,3.2)

#--------------------SLIDES DE REPRESENTADAS------------
def purch_rep_dolar_kg (end_date, start_date, prs, representadas):
    df = df_compras[(df_compras['DT_REL'].dt.month <= (end_date.month))&
                    (df_compras['DT_REL'].dt.month >= (start_date.month))&
                    (df_compras['REPRESENTADA'].isin(representadas))
                     ]
    
    # Process sales data
    chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
    chart_layout_name = 'Grafico_duplo_sem_coment'

    # Find the slide layout by name and add a new slide
    slide_layout = find_layout_by_name(prs, chart_layout_name)

    #-----------------------SLIDE USD x KG----------------------------
    slide = prs.slides.add_slide(slide_layout)

    # Set the title for the slide
    if start_date.month == (end_date.month-2):
        slide.shapes.title.text = f'Purchase - Q{end_date.quarter}/{end_date.year}'
    else:
        slide.shapes.title.text = 'Purchase - YTD'

    #Variáveis
    
    budget = 1
    real = round(df[(df['DT_REL'].dt.year == (end_date.year))]['QTD_NEG'].sum())

    # Add chart 1 - New Projects
    chart_data = CategoryChartData()
    chart_data.categories = ['BUDGET', 'REAL']  # Add the category labels

    # Add the series data for Faturamento and Meta

    chart_data.add_series('KG',(
                        int(budget), 
                        int(real)
                        ))

    #chart_data.add_series('Series 1', (question_data['value']*100).tolist())  # Multiply by 100 to convert to percentage

    # Choose the chart placeholder by index and insert the chart
    chart_placeholder = slide.placeholders[10]  # This is the placeholder index for the chart
    chart_frame = chart_placeholder.insert_chart(chart_type, chart_data)
    chart = chart_frame.chart

    # Formatting based on the chart type
    tratamento_grafico_colunas(chart)

    # Add chart 2 - USD
    real = round(df[(df['DT_REL'].dt.year == (end_date.year))]['VLRTOT_DOLAR'].sum()) 

    chart_data = CategoryChartData()
    chart_data.categories = ['BUDGET', 'REAL']  # Add the category labels

    # Add the series data for Faturamento and Meta

    chart_data.add_series('USD',(
                        int(budget), 
                        int(real)
                        ))
    #chart_data.add_series('Series 1', (question_data['value']*100).tolist())  # Multiply by 100 to convert to percentage

    # Choose the chart placeholder by index and insert the chart
    chart_placeholder = slide.placeholders[13]  # This is the placeholder index for the chart
    chart_frame = chart_placeholder.insert_chart(chart_type, chart_data)
    chart = chart_frame.chart

    # Formatting based on the chart type
    tratamento_grafico_colunas(chart)

def purch_rep_dolar_kg_year (end_date, start_date, prs, representadas):
    df = df_compras[(df_compras['DT_REL'].dt.month <= (end_date.month))&
                    (df_compras['DT_REL'].dt.month >= (start_date.month))&
                    (df_compras['REPRESENTADA'].isin(representadas))
                     ]
    
    # Process sales data
    chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
    chart_layout_name = 'Grafico_duplo_sem_coment'

    # Find the slide layout by name and add a new slide
    slide_layout = find_layout_by_name(prs, chart_layout_name)

    #-----------------------SLIDE USD----------------------------
    slide = prs.slides.add_slide(slide_layout)

    # Set the title for the slide
    if start_date.month == (end_date.month-2):
        slide.shapes.title.text = f'Purchase - Q{end_date.quarter}/{end_date.year} - USD'
    else:
        slide.shapes.title.text = 'Purchase - YTD - USD'
    # Insert text_summary into the specified text placeholder
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame

    #Variáveis
    ano_1 = round(df[(df['DT_REL'].dt.year == (end_date.year -2))]['VLRTOT_DOLAR'].sum())
    ano_2 = round(df[(df['DT_REL'].dt.year == (end_date.year -1))]['VLRTOT_DOLAR'].sum())
    ano_3 = round(df[(df['DT_REL'].dt.year == (end_date.year))]['VLRTOT_DOLAR'].sum())
    budget = 1
    real = ano_3

    # Insert text_summary into the specified text placeholder
    if ano_1 > 0 and ano_2 > 0 and budget > 0:
        text_frame.text = (str(end_date.year-2) + ' x ' + str(end_date.year) + ' = ' + str(round(((ano_3/ano_1)-1)*100)) + '%' + '\n' +
                                str(end_date.year-1) + ' x ' + str(end_date.year) + ' = ' + str(round(((ano_3/ano_2)-1)*100)) + '%' + '\n' +
                                'Budget x Real = ' + str(round(((real/budget)-1)*100)) + '%' + '\n' + 'Validar 2022, pois não há dados completos no Shankhya de compra para esse ano'
                                )
    else:
        text_frame.text = ''

    # Add chart 1 - New Projects
    chart_data = CategoryChartData()
    chart_data.categories = [str(end_date.year-2), str(end_date.year -1), str(end_date.year)]  # Add the category labels

    # Add the series data for Faturamento and Meta

    chart_data.add_series('',(
                        int(ano_1), 
                        int(ano_2),
                        int(ano_3)
                        ))

    #chart_data.add_series('Series 1', (question_data['value']*100).tolist())  # Multiply by 100 to convert to percentage

    # Choose the chart placeholder by index and insert the chart
    chart_placeholder = slide.placeholders[10]  # This is the placeholder index for the chart
    chart_frame = chart_placeholder.insert_chart(chart_type, chart_data)
    chart = chart_frame.chart

    # Formatting based on the chart type
    tratamento_grafico_colunas(chart)

    # Add chart 2 - KG
    chart_data = CategoryChartData()
    chart_data.categories = ['BUDGET', 'REAL']  # Add the category labels

    # Add the series data for Faturamento and Meta

    chart_data.add_series('',(
                        int(budget), 
                        int(real)
                        ))
    #chart_data.add_series('Series 1', (question_data['value']*100).tolist())  # Multiply by 100 to convert to percentage

    # Choose the chart placeholder by index and insert the chart
    chart_placeholder = slide.placeholders[13]  # This is the placeholder index for the chart
    chart_frame = chart_placeholder.insert_chart(chart_type, chart_data)
    chart = chart_frame.chart

    # Formatting based on the chart type
    tratamento_grafico_colunas(chart)

    #---------------------SLIDE KG--------------------------
    slide = prs.slides.add_slide(slide_layout)

    # Set the title for the slide
    if start_date.month == (end_date.month-2):
        slide.shapes.title.text = f'Purchase - Q{end_date.quarter}/{end_date.year} - USD'
    else:
        slide.shapes.title.text = 'Purchase - YTD - USD'

    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame

    #Variáveis
    ano_1 = round(df[(df['DT_REL'].dt.year == (end_date.year -2))]['QTD_NEG'].sum())
    ano_2 = round(df[(df['DT_REL'].dt.year == (end_date.year -1))]['QTD_NEG'].sum())
    ano_3 = round(df[(df['DT_REL'].dt.year == (end_date.year))]['QTD_NEG'].sum())
    budget = 1
    real = ano_3

    # Insert text_summary into the specified text placeholder
    if ano_1 > 0 and ano_2 > 0 and budget > 0:
        text_frame.text = (str(end_date.year-2) + ' x ' + str(end_date.year) + ' = ' + str(round(((ano_3/ano_1)-1)*100)) + '%' + '\n' +
                                str(end_date.year-1) + ' x ' + str(end_date.year) + ' = ' + str(round(((ano_3/ano_2)-1)*100)) + '%' + '\n' +
                                'Budget x Real = ' + str(round(((real/budget)-1)*100)) + '%' + '\n' + 'Validar 2022, pois não há dados completos no Shankhya de compra para esse ano'
                                )
    else:
        text_frame.text = ''

    # Add chart 1 - New Projects
    chart_data = CategoryChartData()
    chart_data.categories = [str(end_date.year-2), str(end_date.year -1), str(end_date.year)]  # Add the category labels

    # Add the series data for Faturamento and Meta

    chart_data.add_series('',(
                        int(ano_1), 
                        int(ano_2),
                        int(ano_3)
                        ))

    #chart_data.add_series('Series 1', (question_data['value']*100).tolist())  # Multiply by 100 to convert to percentage

    # Choose the chart placeholder by index and insert the chart
    chart_placeholder = slide.placeholders[10]  # This is the placeholder index for the chart
    chart_frame = chart_placeholder.insert_chart(chart_type, chart_data)
    chart = chart_frame.chart

    # Formatting based on the chart type
    tratamento_grafico_colunas(chart)

    # Add chart 2 - KG
    chart_data = CategoryChartData()
    chart_data.categories = ['BUDGET', 'REAL']  # Add the category labels

    # Add the series data for Faturamento and Meta

    chart_data.add_series('',(
                        int(budget), 
                        int(real)
                        ))
    #chart_data.add_series('Series 1', (question_data['value']*100).tolist())  # Multiply by 100 to convert to percentage

    # Choose the chart placeholder by index and insert the chart
    chart_placeholder = slide.placeholders[13]  # This is the placeholder index for the chart
    chart_frame = chart_placeholder.insert_chart(chart_type, chart_data)
    chart = chart_frame.chart

    # Formatting based on the chart type
    tratamento_grafico_colunas(chart)

def sales_rep_kg (end_date, start_date, prs, representadas):
    df_sales = df_faturamento[(df_faturamento['Mês/Ano'].dt.month <= end_date.month) &
                            (df_faturamento['Mês/Ano'].dt.month >= start_date.month)&
                            (df_faturamento['REPRESENTADA'].isin(representadas))]
    df_budget = df_meta[(df_meta['Mês/Ano'].dt.month <= end_date.month) &
                            (df_meta['Mês/Ano'].dt.month >= start_date.month)&
                            (df_meta['REPRESENTADA'].isin(representadas))]

    #-----------------MONTAGEM DO SLIDE-------------------------
    # Process sales data

    chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
    chart_layout_name = 'Grafico_duplo_sem_coment'

    # Find the slide layout by name and add a new slide
    slide_layout = find_layout_by_name(prs, chart_layout_name)
    slide = prs.slides.add_slide(slide_layout)

    # Set the title for the slide
    if start_date.month == (end_date.month-2):
        slide.shapes.title.text = f'SALES - Q{end_date.quarter}/{end_date.year}'
    else:
        slide.shapes.title.text = 'SALES - YTD'

    #Variáveis
    ano_1 = df_sales[(df_sales['Mês/Ano'].dt.year == (end_date.year -2))]['KG_FATURAMENTO'].sum()
    ano_2 = df_sales[(df_sales['Mês/Ano'].dt.year == (end_date.year -1))]['KG_FATURAMENTO'].sum()
    ano_3 = df_sales[(df_sales['Mês/Ano'].dt.year == (end_date.year))]['KG_FATURAMENTO'].sum()
    budget = df_budget[(df_budget['Mês/Ano'].dt.year == end_date.year)]['META_KG'].sum()
    real = ano_3

    ano_1_dolar = df_sales[(df_sales['Mês/Ano'].dt.year == (end_date.year -2))]['VLRFAT_TOT_USD'].sum()
    ano_2_dolar = df_sales[(df_sales['Mês/Ano'].dt.year == (end_date.year -1))]['VLRFAT_TOT_USD'].sum()
    ano_3_dolar = df_sales[(df_sales['Mês/Ano'].dt.year == (end_date.year))]['VLRFAT_TOT_USD'].sum()

    budget_dolar = df_budget[(df_budget['Mês/Ano'].dt.year == (end_date.year))]['META_VLR'].sum()
    real_dolar = ano_3_dolar

    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame

    # Insert text_summary into the specified text placeholder
    if ano_1 > 0 and ano_2 > 0 and budget > 0:
        text_frame.text = (str(end_date.year-2) + ' x ' + str(end_date.year) + ' = ' + str(round(((ano_3/ano_1)-1)*100)) + '%' + '\n' +
                                str(end_date.year-1) + ' x ' + str(end_date.year) + ' = ' + str(round(((ano_3/ano_2)-1)*100)) + '%' + '\n' +
                                'Budget x Real = ' + str(round(((real/budget)-1)*100)) + '%' + '\n' +
                                'dolar:'+ '\n' +
                                str(end_date.year-2) + ' x ' + str(end_date.year) + ' = ' + str(round(((ano_3_dolar/ano_1_dolar)-1)*100)) + '%' + '\n' +
                                str(end_date.year-1) + ' x ' + str(end_date.year) + ' = ' + str(round(((ano_3_dolar/ano_2_dolar)-1)*100)) + '%' + '\n' +
                                'Budget x Real = ' + str(round(((real_dolar/budget_dolar)-1)*100)) + '%'
                                )
    else:
        text_frame.text = ''


    # Add chart 1 - New Projects
    chart_data = CategoryChartData()
    chart_data.categories = [str(end_date.year-2), str(end_date.year -1), str(end_date.year)]  # Add the category labels

    # Add the series data for Faturamento and Meta

    chart_data.add_series('KG',(
                        int(ano_1), 
                        int(ano_2),
                        int(ano_3)
                        ))

    #chart_data.add_series('Series 1', (question_data['value']*100).tolist())  # Multiply by 100 to convert to percentage

    # Choose the chart placeholder by index and insert the chart
    chart_placeholder = slide.placeholders[10]  # This is the placeholder index for the chart
    chart_frame = chart_placeholder.insert_chart(chart_type, chart_data)
    chart = chart_frame.chart

    # Formatting based on the chart type
    tratamento_grafico_colunas(chart)

    # Add chart 2 - KG
    chart_data = CategoryChartData()
    chart_data.categories = ['BUDGET', 'REAL']  # Add the category labels

    # Add the series data for Faturamento and Meta

    chart_data.add_series('KG',(
                        int(budget), 
                        int(real)
                        ))
    #chart_data.add_series('Series 1', (question_data['value']*100).tolist())  # Multiply by 100 to convert to percentage

    # Choose the chart placeholder by index and insert the chart
    chart_placeholder = slide.placeholders[13]  # This is the placeholder index for the chart
    chart_frame = chart_placeholder.insert_chart(chart_type, chart_data)
    chart = chart_frame.chart

    # Formatting based on the chart type
    tratamento_grafico_colunas(chart)

def sales_lowlights (end_date, start_date, prs, representadas, produto, lim_table):
    #----------------CRIAÇÃO DO DATAFRAME DE REFERÊNCIA----------------#
    #Oportunidades com previsão para o ano de referência
    df = df_faturamento[
        (df_faturamento['Mês/Ano'].dt.month <= end_date.month) 
        & (df_faturamento['Mês/Ano'].dt.month >= start_date.month) 
        & (df_faturamento['ANO'].isin([end_date.year-1, end_date.year]))
        & (df_faturamento['REPRESENTADA'].isin(representadas))
        ][['ANO','GRUPOECONOMICO','KG_FATURAMENTO']].groupby(['ANO','GRUPOECONOMICO']).sum().reset_index()

    df = df.pivot_table(index='GRUPOECONOMICO', columns='ANO', values='KG_FATURAMENTO').reset_index()
    df = df.fillna(0)
    df['Diferença'] = df[end_date.year] - df[end_date.year-1]

    df_loss = df[(df[end_date.year-1]>0)& (df[end_date.year]>0) & (df['Diferença']<0)].sort_values(by='Diferença', ascending=True)
    #Definindo se vai ou não aparecer o grupo de produto no slide
    if produto == 1:
        df_loss = df_loss.merge(df_faturamento[
            (df_faturamento['ANO']==end_date.year)&
            (df_faturamento['Mês/Ano'].dt.month<=end_date.month)&
            (df_faturamento['Mês/Ano'].dt.month>=start_date.month)&
            (df_faturamento['REPRESENTADA'].isin(representadas))
            ][['GRUPOECONOMICO','DESCRGRUPOPROD']].groupby('GRUPOECONOMICO').agg({'DESCRGRUPOPROD': lambda x : '/ '.join(set(x))}), on='GRUPOECONOMICO', how='left')
        df_loss['GRUPOECONOMICO'] = df_loss['GRUPOECONOMICO'] + ': ' + df_loss['DESCRGRUPOPROD']
    df_loss = df_loss[['GRUPOECONOMICO','Diferença']]
    total_loss = df_loss.shape[0]
    total_loss_kg = df_loss['Diferença'].sum()
    if df_loss.shape[0] > lim_table:
        df_loss = df_loss.sort_values(by='Diferença',ascending = True).head(lim_table)
    df_loss.loc['Total','Diferença'] = total_loss_kg
    df_loss.loc['Total','GRUPOECONOMICO'] = f'Total Loss ({total_loss})'
    df_loss = df_loss.rename(columns={'Diferença':'KG', 'GRUPOECONOMICO':f'Loss vs {end_date.year-1}'})


    df_lost = df[(df[end_date.year-1]>0)& (df[end_date.year]==0)].sort_values(by='Diferença', ascending=True)
    #Definindo se vai ou não aparecer o grupo de produto no slide
    if produto == 1:
        df_lost = df_lost.merge(df_faturamento[
            (df_faturamento['ANO']==end_date.year-1)&
            (df_faturamento['REPRESENTADA'].isin(representadas))]
                                 [['GRUPOECONOMICO','DESCRGRUPOPROD']].groupby('GRUPOECONOMICO').agg({'DESCRGRUPOPROD': lambda x : '/ '.join(set(x))}), on='GRUPOECONOMICO', how='left')
        df_lost['GRUPOECONOMICO'] = df_lost['GRUPOECONOMICO'] + ': ' + df_lost['DESCRGRUPOPROD']
    df_lost = df_lost[['GRUPOECONOMICO','Diferença']]
    total_lost = df_lost.shape[0]
    total_lost_kg = df_lost['Diferença'].sum()
    if df_lost.shape[0] > lim_table:
        df_lost = df_lost.sort_values(by='Diferença',ascending = True).head(lim_table)
    df_lost.loc['Total','Diferença'] = total_lost_kg
    df_lost.loc['Total','GRUPOECONOMICO'] = f'Total Lost ({total_lost})'
    df_lost = df_lost.rename(columns={'Diferença':'KG', 'GRUPOECONOMICO':f'Top Lost Clients vs {end_date.year-1}'})

    df_lost['KG'] = round(df_lost['KG']).astype(int).apply(lambda x: '{:,.0f}'.format(x)).replace(',', '.', regex=True)
    df_loss['KG'] = round(df_loss['KG']).astype(int).apply(lambda x: '{:,.0f}'.format(x)).replace(',', '.', regex=True)
    #---------------CONFIGURAÇÃO DO SLIDE----------------#
    chart_layout_name = 'low_high'

    # Find the slide layout by name and add a new slide
    slide_layout = find_layout_by_name(prs, chart_layout_name)
    slide = prs.slides.add_slide(slide_layout)

    # Set the title for the slide
    if start_date.month == (end_date.month-2):
        slide.shapes.title.text = f'SALES - LOWLIGHTS: Q{end_date.quarter}/{end_date.year}'
    else:
        slide.shapes.title.text = f'SALES - LOWLIGHTS: YTD'

    #---------------TABELA 1 - LOSS ----------------#
    individual_width = [11.5,3.5]
    nova_tabela(df_loss, slide, individual_width, 0.95, 4.20)

    #---------------TABELA 2 - LOST ----------------#
    nova_tabela(df_lost, slide, individual_width, 17.92, 4.20)

def sales_highlights (end_date, start_date, prs, representadas, produto, lim_table):
    #----------------GROWTH----------------#
    df = df_faturamento[(df_faturamento['REPRESENTADA'].isin(representadas))& (df_faturamento['Mês/Ano']<= end_date)]
    df_growth = df[
        (df['Mês/Ano'].dt.month <= end_date.month) &
        (df['Mês/Ano'].dt.month >= start_date.month) & 
        (df['ANO'].isin([end_date.year-1, end_date.year]))
        ][['ANO','GRUPOECONOMICO','KG_FATURAMENTO']].groupby(['ANO','GRUPOECONOMICO']).sum().reset_index()

    df_growth = df_growth.pivot_table(index='GRUPOECONOMICO', columns='ANO', values='KG_FATURAMENTO').reset_index()
    df_growth = df_growth.fillna(0)
    df_growth['Diferença'] = df_growth[end_date.year] - df_growth[end_date.year-1]

    df_growth = df_growth[(df_growth[end_date.year-1]>0)& (df_growth[end_date.year]>0) & (df_growth['Diferença']>0)].sort_values(by='Diferença', ascending=False)
    if produto == 1:
        df_growth = df_growth.merge(df[
            (df['ANO']==end_date.year)&
            (df['Mês/Ano'].dt.month<=end_date.month)&
            (df['Mês/Ano'].dt.month>=start_date.month)&
            (df['REPRESENTADA'].isin(representadas))
                                    ][['GRUPOECONOMICO','DESCRGRUPOPROD']].groupby('GRUPOECONOMICO').agg({'DESCRGRUPOPROD': lambda x : '/ '.join(set(x))}), on='GRUPOECONOMICO', how='left')
        df_growth['GRUPOECONOMICO'] = df_growth['GRUPOECONOMICO'] + ': ' + df_growth['DESCRGRUPOPROD']
    df_growth = df_growth[['GRUPOECONOMICO','Diferença']]
    total_growth = df_growth.shape[0]
    total_growth_kg = df_growth['Diferença'].sum()
    if df_growth.shape[0] > lim_table:
        df_growth = df_growth.sort_values(by='Diferença',ascending = False).head(lim_table)
    df_growth.loc['Total','Diferença'] = total_growth_kg
    df_growth.loc['Total','GRUPOECONOMICO'] = f'Total Loss ({total_growth})'
    df_growth = df_growth.rename(columns={'Diferença':'KG', 'GRUPOECONOMICO':f'Growth vs {end_date.year-1}'})
    df_growth['KG'] = round(df_growth['KG']).astype(int)

    #----------------REACTIVATIONS----------------#
    #Regra: clientes que voltaram a faturar depois de 12 meses sem faturamento
    df_react = df[['Mês/Ano','GRUPOECONOMICO','KG_FATURAMENTO']].groupby(['Mês/Ano','GRUPOECONOMICO']).sum().reset_index()
    df_react['FIRST_FAT_YEAR'] = df_react[df_react['Mês/Ano'].dt.year == end_date.year].groupby('GRUPOECONOMICO')['Mês/Ano'].transform('min')
    df_react['LAST_PREV'] = df_react[df_react['Mês/Ano'].dt.year != end_date.year].groupby('GRUPOECONOMICO')['Mês/Ano'].transform('max')
    df_react['FAT_CURRENT_YEAR'] = df_react[df_react['Mês/Ano'].dt.year == end_date.year].groupby('GRUPOECONOMICO')['KG_FATURAMENTO'].transform('sum')
    df_react = df_react[['GRUPOECONOMICO','FIRST_FAT_YEAR','LAST_PREV','KG_FATURAMENTO','FAT_CURRENT_YEAR']].groupby('GRUPOECONOMICO').agg({'FIRST_FAT_YEAR':'min','LAST_PREV':'max','FAT_CURRENT_YEAR':'sum'}).reset_index()
    df_react = df_react[(df_react['LAST_PREV'].notnull()) & (df_react['FIRST_FAT_YEAR'].notnull())]
    df_react = df_react[(df_react['FIRST_FAT_YEAR'].dt.month >= start_date.month)]
    df_react['LAST_PREV'] = pd.to_datetime(df_react['LAST_PREV'])
    df_react['FIRST_FAT_YEAR'] = pd.to_datetime(df_react['FIRST_FAT_YEAR'])
    df_react['DIFF'] = df_react['FIRST_FAT_YEAR'] - df_react['LAST_PREV']
    df_react['DIFF'] = df_react['DIFF'].dt.days
    df_react = df_react[df_react['DIFF'] > 365].sort_values(by='FAT_CURRENT_YEAR', ascending=False)
    if produto == 1:
        df_react = df_react.merge(df[(df['ANO']==end_date.year)][['GRUPOECONOMICO','DESCRGRUPOPROD']].groupby('GRUPOECONOMICO').agg({'DESCRGRUPOPROD': lambda x : '/ '.join(set(x))}), on='GRUPOECONOMICO', how='left')
        df_react['GRUPOECONOMICO'] = df_react['GRUPOECONOMICO'] + ' - ' + df_react['DESCRGRUPOPROD']
    df_react = df_react.rename(columns={'GRUPOECONOMICO':'Reactivation - Sales Kg','FAT_CURRENT_YEAR':'KG'})
    df_react = df_react[['Reactivation - Sales Kg','KG']]
    total_react = df_react.shape[0]
    total_react_kg = df_react['KG'].sum()
    if df_react.shape[0] > lim_table:
        df_react = df_react.sort_values(by='KG',ascending = False).head(lim_table)
    df_react.loc['Total','KG'] = total_react_kg
    df_react.loc['Total','Reactivation - Sales Kg'] = f'Total Reactivations ({total_react})'

    #----------------NEW CLIENTS----------------#
    df_new = df[['Mês/Ano','GRUPOECONOMICO','KG_FATURAMENTO']].groupby(['GRUPOECONOMICO']).agg({'Mês/Ano':'min','KG_FATURAMENTO':'sum'}).reset_index()
    df_new['Mês/Ano'] = pd.to_datetime(df_new['Mês/Ano'])
    df_new = df_new[(df_new['Mês/Ano'] > start_date)]
    if produto == 1:
        df_new = df_new.merge(df[df['ANO']== end_date.year][['GRUPOECONOMICO','DESCRGRUPOPROD']].groupby('GRUPOECONOMICO').agg({'DESCRGRUPOPROD':lambda x : '/'.join(set(x))}), on = 'GRUPOECONOMICO', how='left')
        df_new['GRUPOECONOMICO'] = df_new['GRUPOECONOMICO'] + ' - ' + df_new['DESCRGRUPOPROD']
    df_new = df_new.rename(columns={'KG_FATURAMENTO':'KG','GRUPOECONOMICO':'New Clients - Sales Kg'})
    df_new = df_new[['New Clients - Sales Kg','KG']]
    total_new = df_new.shape[0]
    total_new_kg = df_new['KG'].sum()
    if df_new.shape[0] > lim_table:
        df_new = df_new.sort_values(by='KG',ascending = False).head(lim_table)
    df_new.loc['Total','KG'] = total_new_kg
    df_new.loc['Total','New Clients - Sales Kg'] = f'Total New ({total_new})'

    #Definição de formatação de KG
    df_growth['KG'] = df_growth['KG'].apply(lambda x: '{:,.0f}'.format(x)).replace(',', '.', regex=True)
    df_react['KG'] = df_react['KG'].apply(lambda x: '{:,.0f}'.format(x)).replace(',', '.', regex=True)
    df_new['KG'] = df_new['KG'].apply(lambda x: '{:,.0f}'.format(x)).replace(',', '.', regex=True)

    #---------------CONFIGURAÇÃO DO SLIDE----------------#
    chart_layout_name = 'low_high'

    # Find the slide layout by name and add a new slide
    slide_layout = find_layout_by_name(prs, chart_layout_name)
    slide = prs.slides.add_slide(slide_layout)

    # Set the title for the slide
    if start_date.month == (end_date.month-2):
        slide.shapes.title.text = f'SALES - HIGHILIGHTS: Q{end_date.quarter}/{end_date.year}'
    else:
        slide.shapes.title.text = f'SALES - HIGHILIGHTS: YTD'

    #---------------TABELA 1 - GROWTH ----------------#
    individual_width = [6.82,3.08]
    nova_tabela(df_growth, slide, individual_width, 0.77, 3.80)

    #---------------TABELA 2 - REACTIVATIONS ----------------#
    nova_tabela(df_react, slide, individual_width, 11.98, 3.80)

    #---------------TABELA 3 - NEW CLIENTS ----------------#
    nova_tabela(df_new, slide, individual_width, 23.2, 3.80)

def oport_highlights (end_date, start_date, prs, representadas, lbz, bu):

    df = df_oportunidade[(df_oportunidade['REPRESENTADA'].isin(representadas)) & (df_oportunidade['DTNEG'].dt.month >= start_date.month) & (df_oportunidade['DTNEG'].dt.month <= end_date.month)]

    if lbz == 1:
        df = df[df['DESCRICAO_APLICACAO'].str.contains(bu[0])|df['DESCRICAO_APLICACAO'].str.contains(bu[1])]
    
    chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
    chart_layout_name = 'Grafico_duplo'

    # Find the slide layout by name and add a new slide
    slide_layout = find_layout_by_name(prs, chart_layout_name)
    slide = prs.slides.add_slide(slide_layout)

    # Set the title for the slide
    slide.shapes.title.text = 'PROJECTS HIGHLIGHTS'

    # Insert text_summary into the specified text placeholder
    text_placeholder = slide.placeholders[12]  # Index 11 for the text_summary placeholder
    text_placeholder.text = 'New projects ' + str((start_date.year - 1)) + 'x' + str((start_date.year)) + ' Between ' + str(start_date.month) + ' and ' + str(end_date.month) + (' - ' + bu[0] + ' / ' + bu[1] if lbz == 1 else '')

    # Add chart 1 - New Projects
    chart_data = CategoryChartData()
    chart_data.categories = ['2023', '2024']  # Add the category labels

    # Add the series data for Faturamento and Meta

    chart_data.add_series('New Projects',(
                        int(df[(df['DTNEG'].dt.year == (end_date.year -1))&(df['STATUS'] != 'Canceled')]['NUNEGOCIACAO'].count()), 
                        int(df[(df['DTNEG'].dt.year == end_date.year)&(df['STATUS'] != 'Canceled')]['NUNEGOCIACAO'].count())))

    #chart_data.add_series('Series 1', (question_data['value']*100).tolist())  # Multiply by 100 to convert to percentage

    # Choose the chart placeholder by index and insert the chart
    chart_placeholder = slide.placeholders[10]  # This is the placeholder index for the chart
    chart_frame = chart_placeholder.insert_chart(chart_type, chart_data)
    chart = chart_frame.chart

    # Formatting based on the chart type
    tratamento_grafico_colunas(chart)

    # Add chart 2 - KG
    chart_data = CategoryChartData()
    chart_data.categories = ['2023', '2024']  # Add the category labels

    # Add the series data for Faturamento and Meta

    chart_data.add_series('KG',(
                        int(df[(df['DTNEG'].dt.year == (end_date.year -1))&(df['STATUS'] != 'Canceled')]['QTDNEG'].sum()), 
                        int(df[(df['DTNEG'].dt.year == end_date.year)&(df['STATUS'] != 'Canceled')]['QTDNEG'].sum())))

    #chart_data.add_series('Series 1', (question_data['value']*100).tolist())  # Multiply by 100 to convert to percentage

    # Choose the chart placeholder by index and insert the chart
    chart_placeholder = slide.placeholders[13]  # This is the placeholder index for the chart
    chart_frame = chart_placeholder.insert_chart(chart_type, chart_data)
    chart = chart_frame.chart

    tratamento_grafico_colunas(chart)

def oport_em_aberto (end_date, start_date, prs, representadas, ano_vigente, lbz, bu, lim_table):
    if lbz == 1:
        df = df_oportunidade[df_oportunidade['DESCRICAO_APLICACAO'].str.contains(bu[0])|df_oportunidade['DESCRICAO_APLICACAO'].str.contains(bu[1])]
    else:
        df = df_oportunidade

    df = df[
        (df['STATUS RESUMIDO'].isin(['Em Andamento'])) &
        (df['REPRESENTADA'].isin(representadas)) &
        (df['DTNEG'] <= (end_date + pd.DateOffset(months=1)))
        ][['GRUPO_ECONOMICO','GRUPO_PRODUTO','NOME_PROJETO','DTESTFECHAMENTO','VLRTOT','QTDNEG']].sort_values(by='QTDNEG', ascending=False)

    if ano_vigente == 1:
        df = df[(df['DTESTFECHAMENTO'].dt.year == end_date.year)]

    total_projects = str(df.shape[0])

    df = df[['GRUPO_ECONOMICO','GRUPO_PRODUTO','NOME_PROJETO','QTDNEG','VLRTOT','DTESTFECHAMENTO']]
    df = df.merge(tabela_compras, on = 'GRUPO_PRODUTO', how='left')
    df = df.fillna(0)
    df['VLRTOT'] = df['QTDNEG'] * df['VLRUNITDOLAR']
    df = df[['GRUPO_ECONOMICO','GRUPO_PRODUTO','NOME_PROJETO','QTDNEG','VLRTOT','DTESTFECHAMENTO']].sort_values(by='QTDNEG',ascending = False)
    df.rename(columns={'GRUPO_ECONOMICO': 'CUSTOMER', 'NOME_PROJETO': 'PROJECT', 'GRUPO_PRODUTO': 'PRODUCT', 'QTDNEG': 'ANUAL VOLUME (KG)', 'VLRTOT': 'VALUE (USD)', 'DTESTFECHAMENTO': 'CONVERSION ESTIMATED DATE'}, inplace=True)
    df['CONVERSION ESTIMATED DATE'] = df['CONVERSION ESTIMATED DATE'].dt.strftime('%d/%m/%Y')
    peso_total = df['ANUAL VOLUME (KG)'].sum() 
    dolar_total = df['VALUE (USD)'].sum()
    if df.shape[0] > lim_table:
        df = df.head(lim_table)
    if df.shape[0] > 0:
        df.loc['Total','ANUAL VOLUME (KG)'] = peso_total
        df.loc['Total','VALUE (USD)'] = dolar_total
        df.loc['Total','CUSTOMER'] =f'Total Projects: {str(total_projects)}'
        df['ANUAL VOLUME (KG)'] = df['ANUAL VOLUME (KG)'].apply(lambda x: '{:,.0f}'.format(x)).str.replace(',', '.')
        df['VALUE (USD)'] = df['VALUE (USD)'].apply(lambda x: '$' + '{:,.0f}'.format(x)).str.replace(',', '.')
    peso_total = '{:,.0f}'.format(peso_total).replace(',', '.')

    #---------------CONFIGURAÇÃO DO SLIDE----------------#
    #---------------GRÁFICO----------------#

    chart_type = XL_CHART_TYPE.PIE
    chart_layout_name = 'Tabela_Grafico_Sem_Coment'

    # Find the slide layout by name and add a new slide
    slide_layout = find_layout_by_name(prs, chart_layout_name)
    slide = prs.slides.add_slide(slide_layout)

    # Set the title for the slide
    slide.shapes.title.text = 'HIGHLIGHTS'

    # Insert text_summary into the specified text placeholder
    text_placeholder = slide.placeholders[14]  # Index 11 for the text_summary placeholder
    if ano_vigente == 1:
        text_placeholder.text = 'Estimated for ' + str((start_date.year)) + (' - ' + bu[0] + ' / ' + bu[1] if lbz == 1 else '')
    else:
        text_placeholder.text = 'Open Projects' + (' - ' + bu[0] + ' / ' + bu[1] if lbz == 1 else '')

    # Add chart 1 - New Projects
    chart_data = CategoryChartData()
    chart_data.categories = ['Approved', 'First Steps','Stability','Formulations Test','Final Tests', 'Negociation'] # Add the category labels

    # Add the series data for Faturamento and Meta

    if ano_vigente == 1:
        df_chart = df_oportunidade[(df_oportunidade['DTESTFECHAMENTO'].dt.year == end_date.year)&(df_oportunidade['REPRESENTADA'].isin(representadas))]
    else:
        df_chart = df_oportunidade[(df_oportunidade['REPRESENTADA'].isin(representadas))]
    total_oportunidades_pipe = int(df_chart[(df_chart['STATUS'].isin(['Approved', 'First Steps', 'Stability', 'Formulations Test', 'Final Tests', 'Negociation']))]['NUNEGOCIACAO'].count())
    if total_oportunidades_pipe == 0:
        total_oportunidades_pipe = 1

    chart_data.add_series('STATUS (%)', [
        int(df_chart[(df_chart['STATUS'] == 'Approved')]['NUNEGOCIACAO'].count()/ total_oportunidades_pipe*100),
        int(df_chart[(df_chart['STATUS'] == 'First Steps')]['NUNEGOCIACAO'].count()/ total_oportunidades_pipe*100),
        int(df_chart[(df_chart['STATUS'] == 'Stability')]['NUNEGOCIACAO'].count()/ total_oportunidades_pipe*100),
        int(df_chart[(df_chart['STATUS'] == 'Formulations Test')]['NUNEGOCIACAO'].count()/ total_oportunidades_pipe*100),
        int(df_chart[(df_chart['STATUS'] == 'Final Tests')]['NUNEGOCIACAO'].count()/ total_oportunidades_pipe*100),
        int(df_chart[(df_chart['STATUS'] == 'Negociation')]['NUNEGOCIACAO'].count()/ total_oportunidades_pipe*100)
    ])

    chart_placeholder = slide.placeholders[13]  # This is the placeholder index for the chart
    chart_frame = chart_placeholder.insert_chart(chart_type, chart_data)
    chart = chart_frame.chart

    # Formatting based on the chart type
    for series in chart.series:
        fill = series.format.fill
        fill.solid()

        # Add data labels formatted as percentages with no decimal points
        series.has_data_labels = True
        for point in series.points:
            point.data_label.number_format = '#.##0'
            point.data_label.font.size = Pt(14)
            point.data_label.font.typeface = 'Futura Lt BT'
            point.data_label.font.color.rgb = RGBColor(0x00, 0x00, 0x00)  # Black font color
        chart.has_title = True 
        chart.has_legend = True
        chart.legend.font.size = Pt(11)
        chart.legend.font.name = 'Futura Lt BT'
        chart.legend.include_in_layout = False
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM  

    #---------------TABELA----------------#

    individual_width = [3.53,3.61,5.25,2.47,2.55,3.76]
    nova_tabela(df, slide, individual_width,0.22,3.14)

def oport_abertas (end_date, start_date, prs, representadas, lbz, bu, lim_table):
    #----------------CRIAÇÃO DO DATAFRAME DE REFERÊNCIA----------------#
    #Oportunidades com previsão para o ano de referência
    if lbz == 1:
        df = df_oportunidade[df_oportunidade['DESCRICAO_APLICACAO'].str.contains(bu[0])|df_oportunidade['DESCRICAO_APLICACAO'].str.contains(bu[1])]
    else:
        df = df_oportunidade

    df = df[(df['DTNEG'].dt.month <= end_date.month) 
                                      & (df['DTNEG'].dt.month >= start_date.month) 
                                      & (df['DTNEG'].dt.year == end_date.year) 
                                      & (df['STATUS RESUMIDO'] != 'Cancelado') 
                                      & (df['REPRESENTADA'].isin(representadas))
                                      ][['GRUPO_ECONOMICO','GRUPO_PRODUTO','NOME_PROJETO','DTESTFECHAMENTO','QTDNEG','VLRTOT','STATUS']].sort_values(by='QTDNEG', ascending=False)

    total_projects = str(df.shape[0])

    df = df.merge(tabela_compras, on = 'GRUPO_PRODUTO', how='left')
    df = df.fillna(0)
    df['VLRTOT'] = df['QTDNEG'] * df['VLRUNITDOLAR']
    df = df[['GRUPO_ECONOMICO','NOME_PROJETO','GRUPO_PRODUTO','STATUS','QTDNEG','VLRTOT','DTESTFECHAMENTO']]
    df.rename(columns={'VLRTOT':'VALUE (USD)','GRUPO_ECONOMICO': 'CUSTOMER', 'NOME_PROJETO': 'PROJECT', 'GRUPO_PRODUTO': 'PRODUCT', 'QTDNEG': 'ANUAL VOLUME (KG)', 'DTESTFECHAMENTO': 'CONVERSION ESTIMATED DATE'}, inplace=True)
    df['CONVERSION ESTIMATED DATE'] = df['CONVERSION ESTIMATED DATE'].dt.strftime('%d/%m/%Y')
    peso_total = df['ANUAL VOLUME (KG)'].sum() 
    dolar_total = df['VALUE (USD)'].sum()
    total_projects = len(df)
    if df.shape[0] > lim_table:
        df = df.head(lim_table)
    if df.shape[0] > 0:
        df.loc['Total', 'ANUAL VOLUME (KG)'] = peso_total
        df.loc['Total', 'VALUE (USD)'] = dolar_total
        df.loc['Total', 'PROJECT'] = f'{str(total_projects)} Projects'
        df.loc['Total', 'CUSTOMER'] = 'Total'
        df['ANUAL VOLUME (KG)'] = df['ANUAL VOLUME (KG)'].apply(lambda x: '{:,.0f}'.format(x)).str.replace(',', '.')
        df['VALUE (USD)'] = df['VALUE (USD)'].apply(lambda x: '$' + '{:,.0f}'.format(x)).str.replace(',', '.')

    #---------------CONFIGURAÇÃO DO SLIDE----------------#
    chart_layout_name = 'tabela_simples_sem_coment'

    # Find the slide layout by name and add a new slide
    slide_layout = find_layout_by_name(prs, chart_layout_name)
    slide = prs.slides.add_slide(slide_layout)

    # Set the title for the slide
    slide.shapes.title.text = 'PROJECTS - OPENED PROJECTS'

    # Insert text_summary into the specified text placeholder
    text_placeholder = slide.placeholders[15]  # Index 11 for the text_summary placeholder
    text_placeholder.text = 'New projects Between ' + str(start_date.month) + ' and ' + str(end_date.month) + '/' + str(end_date.year) + (' - ' + bu[0] + ' / ' + bu[1] if lbz == 1 else '')

    #---------------TABELA----------------#
    individual_width = [5.08,9.5,5.66,3.2,2.71,2.75,3.83]
    nova_tabela(df, slide, individual_width,0.57,3.01)

def oport_perdidas_chart_table_reproved (end_date, start_date, prs, representadas, lim_table, lbz, bu):
    if lbz == 1:
        df = df_oportunidade[df_oportunidade['DESCRICAO_APLICACAO'].str.contains(bu[0])|df_oportunidade['DESCRICAO_APLICACAO'].str.contains(bu[1])]
    else:
        df = df_oportunidade

    df = df[(df['DTFECHAMENTO'].dt.year == end_date.year) &
                         (df['DTFECHAMENTO'].dt.month <= end_date.month) & 
                         (df['DTFECHAMENTO'].dt.month >= start_date.month) &
                         (df['REPRESENTADA'].isin(representadas)) & 
                         (df['STATUS'].isin(['Reproved','Canceled']))
    ][['NUNEGOCIACAO','GRUPO_ECONOMICO','GRUPO_PRODUTO','QTDNEG','VLRTOT','STATUS','MOTIVO_STATUS']]

    #--------------SLIDES-------------------
    chart_type = XL_CHART_TYPE.PIE
    chart_layout_name = 'Grafico_simples'

    # Find the slide layout by name and add a new slide
    slide_layout = find_layout_by_name(prs, chart_layout_name)
    slide = prs.slides.add_slide(slide_layout)

    # Set the title for the slide
    slide.shapes.title.text = 'PPROJECTS - LOST - BETWEEN - ' + str(start_date.month) + ' AND ' + str(end_date.month) + '/' + str(end_date.year) + (' - ' + bu[0] + ' / ' + bu[1] if lbz == 1 else '')

    # Insert text_summary into the specified text placeholder
    text_placeholder = slide.placeholders[12]  # Index 11 for the text_summary placeholder
    text_placeholder.text = 'Total projects lost: ' + str(df.shape[0]) 

    # Add chart 1 - Personal Care
    chart_data = CategoryChartData()
    chart_data.categories = ['Canceled', 'Reproved'] # Add the category labels

    # Add the series data for Faturamento and Meta

    total_oportunidades_pipe = df.shape[0]

    chart_data.add_series('(%)', [
        int(df[(df['STATUS']=='Canceled')]['NUNEGOCIACAO'].count()/ total_oportunidades_pipe*100),
        int(df[(df['STATUS']=='Reproved')]['NUNEGOCIACAO'].count()/ total_oportunidades_pipe*100),
    ])

    chart_placeholder = slide.placeholders[10]  # This is the placeholder index for the chart
    chart_frame = chart_placeholder.insert_chart(chart_type, chart_data)
    chart = chart_frame.chart

    # Formatting based on the chart type
    for series in chart.series:
        fill = series.format.fill
        fill.solid()

        # Add data labels formatted as percentages with no decimal points
        series.has_data_labels = True
        for point in series.points:
            point.data_label.number_format = '#.##0'
            point.data_label.font.size = Pt(14)
            point.data_label.font.typeface = 'Futura Lt BT'
            point.data_label.font.color.rgb = RGBColor(0x00, 0x00, 0x00)  # Black font color
        chart.has_title = True 
        chart.has_legend = True
        chart.legend.font.size = Pt(11)
        chart.legend.font.name = 'Futura Lt BT'
        chart.legend.include_in_layout = False
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM

    #---------------SLIDE REPROVADOS--------------------
    df = df[df['STATUS'] == 'Reproved'].sort_values(by='QTDNEG', ascending=False)
    total_projects = str(df.shape[0])

    df = df.merge(tabela_compras, on = 'GRUPO_PRODUTO', how='left')
    df = df.fillna(0)
    df['VLRTOT'] = df['QTDNEG'] * df['VLRUNITDOLAR']
    df = df[['GRUPO_ECONOMICO','GRUPO_PRODUTO','QTDNEG','VLRTOT','MOTIVO_STATUS']]
    df.rename(columns={'GRUPO_ECONOMICO': 'CUSTOMER', 'GRUPO_PRODUTO': 'PRODUCT', 'QTDNEG': 'ANUAL VOLUME (KG)', 'VLRTOT':'VALUE (USD)', 'MOTIVO_STATUS':'COMMENTS'}, inplace=True)
    peso_total = df['ANUAL VOLUME (KG)'].sum()
    dolar_total = df['VALUE (USD)'].sum()
    if df.shape[0] > lim_table:
        df = df.head(lim_table)
    if df.shape[0] > 0:
        df.loc['Total','ANUAL VOLUME (KG)'] = peso_total
        df.loc['Total','VALUE (USD)'] = dolar_total
        df.loc['Total','CUSTOMER'] = f'Total Projects: {str(total_projects)}'
        df['ANUAL VOLUME (KG)'] = df['ANUAL VOLUME (KG)'].apply(lambda x: '{:,.0f}'.format(x)).str.replace(',', '.')
        df['VALUE (USD)'] = df['VALUE (USD)'].apply(lambda x: '$' + '{:,.0f}'.format(x)).str.replace(',', '.')
    peso_total = '{:,.0f}'.format(peso_total).replace(',', '.')
    dolar_total = '{:,.0f}'.format(dolar_total).replace(',', '.')

    #---------------CONFIGURAÇÃO DO SLIDE----------------#
    chart_layout_name = 'tabela_simples_sem_coment'

    # Find the slide layout by name and add a new slide
    slide_layout = find_layout_by_name(prs, chart_layout_name)
    slide = prs.slides.add_slide(slide_layout)

    # Set the title for the slide
    slide.shapes.title.text = 'PROJECTS - REPROVED'

    # Insert text_summary into the specified text placeholder
    text_placeholder = slide.placeholders[15]  # Index 11 for the text_summary placeholder
    text_placeholder.text = 'Between ' + str(start_date.month) + ' and ' + str(end_date.month) + '/' + str(end_date.year)

    #---------------TABELA----------------#
    individual_width = [9.8,9.8,3.55,2.45,7.65]
    nova_tabela(df, slide, individual_width,0.3,3.24)

def oport_convertidas (end_date, start_date, prs, representadas, lbz, bu, lim_table):
    if lbz == 1:
        df = df_faturamento[df_faturamento['SEG_PRINCIPAL'].str.contains(bu[0])|df_faturamento['SEG_PRINCIPAL'].str.contains(bu[1])]
    else:
        df = df_faturamento
    df = df[(df['TIPMOV'] == 'V-Venda')&(df['REPRESENTADA'].isin(representadas))][['Mês/Ano','GRUPOECONOMICO','DESCRGRUPOPROD','KG_FATURAMENTO']].groupby(['GRUPOECONOMICO', 'DESCRGRUPOPROD']).agg({'Mês/Ano':'min','KG_FATURAMENTO':'sum'}).reset_index()
    df['Mês/Ano'] = pd.to_datetime(df['Mês/Ano'])
    df = df[(df['Mês/Ano'].dt.year == start_date.year)& (df['Mês/Ano'].dt.month >= start_date.month)&(df['Mês/Ano'].dt.month <= end_date.month)]
    df = df[['GRUPOECONOMICO','DESCRGRUPOPROD','KG_FATURAMENTO']]
    df = df.merge(tabela_compras, left_on = 'DESCRGRUPOPROD', right_on = 'GRUPO_PRODUTO', how='left')
    df = df.fillna(0)
    df['VLRTOT'] = df['KG_FATURAMENTO'] * df['VLRUNITDOLAR']
    df = df[['GRUPOECONOMICO','DESCRGRUPOPROD','KG_FATURAMENTO','VLRTOT']]
    df = df.rename(columns={'KG_FATURAMENTO':'SALES (KG)','DESCRGRUPOPROD':'PRODUCT','GRUPOECONOMICO':'CUSTOMER','VLRTOT':'VALUE (USD)'})

    total_projects = df.shape[0]
    peso_total = df['SALES (KG)'].sum()
    US_total = df['VALUE (USD)'].sum()

    if df.shape[0] > lim_table:
        df = df.sort_values(by='SALES (KG)', ascending= False).head(lim_table)
    if df.shape[0] > 0:
        df.loc['Total','SALES (KG)'] = peso_total
        df.loc['Total','VALUE (USD)'] = US_total
        df.loc['Total','CUSTOMER'] = f'Total Projects ({total_projects})'
        df['SALES (KG)'] = df['SALES (KG)'].apply(lambda x: '{:,.0f}'.format(x)).str.replace(',', '.')
        df['VALUE (USD)'] = df['VALUE (USD)'].apply(lambda x: '$' + '{:,.0f}'.format(x)).str.replace(',', '.')
    US_total = '{:,.0f}'.format(US_total).replace(',', '.')
    peso_total = '{:,.0f}'.format(peso_total).replace(',', '.')

    #---------------CONFIGURAÇÃO DO SLIDE----------------#
    chart_layout_name = 'tabela_simples_sem_coment'

    # Find the slide layout by name and add a new slide
    slide_layout = find_layout_by_name(prs, chart_layout_name)
    slide = prs.slides.add_slide(slide_layout)

    # Set the title for the slide
    slide.shapes.title.text = 'PROJECTS - CONVERTED'

    # Insert text_summary into the specified text placeholder
    text_placeholder = slide.placeholders[15]  # Index 11 for the text_summary placeholder
    text_placeholder.text = str(start_date.month) + '-' + str(end_date.month) + '/' + str(end_date.year) + (' - ' + bu[0] + ' / ' + bu[1] if lbz == 1 else '')

    #---------------TABELA----------------#
    individual_width = [12.68,13.28,3.48,3.67]
    nova_tabela(df, slide, individual_width,0.31,3.26)

def oport_convertida_por_quarter (start_date, prs, representadas, lbz, bu):
    if lbz == 1:
        df = df_faturamento[df_faturamento['SEG_PRINCIPAL'].str.contains(bu[0])|df_faturamento['SEG_PRINCIPAL'].str.contains(bu[1])]
    else:
        df = df_faturamento
    df = df[(df['TIPMOV'] == 'V-Venda')&(df['REPRESENTADA'].isin(representadas))][['Mês/Ano','GRUPOECONOMICO','DESCRGRUPOPROD','KG_FATURAMENTO']].groupby(['GRUPOECONOMICO', 'DESCRGRUPOPROD']).agg({'Mês/Ano':'min','KG_FATURAMENTO':'sum'}).reset_index()
    df['Mês/Ano'] = pd.to_datetime(df['Mês/Ano'])
    df = df[(df['Mês/Ano'].dt.year == start_date.year)]
    df['QUARTER'] = df['Mês/Ano'].dt.quarter
    df['FATURAMENTO_QUARTER'] = df.apply(lambda row: df_faturamento[(df_faturamento['GRUPOECONOMICO'] == row['GRUPOECONOMICO']) & 
                                                                    (df_faturamento['DESCRGRUPOPROD'] == row['DESCRGRUPOPROD']) & 
                                                                    (df_faturamento['Mês/Ano'].dt.quarter == row['QUARTER'])]['KG_FATURAMENTO'].sum(), axis=1)
    df = df[['GRUPOECONOMICO','DESCRGRUPOPROD','QUARTER', 'FATURAMENTO_QUARTER', 'KG_FATURAMENTO']]
    df = df.rename(columns={'FATURAMENTO_QUARTER':'SALES (KG)','DESCRGRUPOPROD':'PRODUCT','GRUPOECONOMICO':'CUSTOMER','QUARTER':'QUARTER'})

    #---------------CONFIGURAÇÃO DO SLIDE----------------#
    chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
    chart_layout_name = 'Grafico_duplo_sem_coment'

    # Find the slide layout by name and add a new slide
    slide_layout = find_layout_by_name(prs, chart_layout_name)

    #-----------------------SLIDE USD x KG----------------------------
    slide = prs.slides.add_slide(slide_layout)

    # Set the title for the slide
    
    slide.shapes.title.text = 'CONVERTED PROJECTS ' + str(start_date.year) + (' - ' + bu[0] + ' / ' + bu[1] if lbz == 1 else '')

    #Variáveis

    # Add chart 1 - New Projects
    chart_data = CategoryChartData()
    chart_data.categories = ['Q1', 'Q2', 'Q3', 'Q4']  # Add the category labels

    # Add the series data for Faturamento and Meta

    chart_data.add_series('Number of Projects',(df[df['QUARTER'] == 1]['CUSTOMER'].count(),
                                            df[df['QUARTER'] == 2]['CUSTOMER'].count(),
                                            df[df['QUARTER'] == 3]['CUSTOMER'].count(),
                                            df[df['QUARTER'] == 4]['CUSTOMER'].count()))

    #chart_data.add_series('Series 1', (question_data['value']*100).tolist())  # Multiply by 100 to convert to percentage

    # Choose the chart placeholder by index and insert the chart
    chart_placeholder = slide.placeholders[10]  # This is the placeholder index for the chart
    chart_frame = chart_placeholder.insert_chart(chart_type, chart_data)
    chart = chart_frame.chart

    # Formatting based on the chart type
    tratamento_grafico_colunas(chart)

    # Add chart 2 - USD
    chart_data = CategoryChartData()
    chart_data.categories = ['Q1', 'Q2', 'Q3', 'Q4']  # Add the category labels

    # Add the series data for Faturamento and Meta

    chart_data.add_series('KG',(df[df['QUARTER'] == 1]['KG_FATURAMENTO'].sum(),
                                            df[df['QUARTER'] == 2]['KG_FATURAMENTO'].sum(),
                                            df[df['QUARTER'] == 3]['KG_FATURAMENTO'].sum(),
                                            df[df['QUARTER'] == 4]['KG_FATURAMENTO'].sum()))
    #chart_data.add_series('Series 1', (question_data['value']*100).tolist())  # Multiply by 100 to convert to percentage

    # Choose the chart placeholder by index and insert the chart
    chart_placeholder = slide.placeholders[13]  # This is the placeholder index for the chart
    chart_frame = chart_placeholder.insert_chart(chart_type, chart_data)
    chart = chart_frame.chart

    # Formatting based on the chart type
    tratamento_grafico_colunas(chart)

#------------------DEMANDAS SPOT--------------------
def oport_grupo_especifico (prs, group):
    #----------------CRIAÇÃO DO DATAFRAME DE REFERÊNCIA----------------#
    #Oportunidades com previsão para o ano de referência
    df = df_oportunidade[
        (df_oportunidade['DESCRICAO_APLICACAO'].str.contains('I&I')|df_oportunidade['DESCRICAO_APLICACAO'].str.contains('HOUSEHOLD')) &
        (df_oportunidade['STATUS RESUMIDO'] == 'Em Andamento') &
        (df_oportunidade['GRUPO_PRODUTO'].isin(['CARBOPOL EZ 4','NOVERITE 311 POLYMER','NOVERITE LD920N','DILUTHIX CLEAR SOFTENER']))
        ][['GRUPO_ECONOMICO','GRUPO_PRODUTO','NOME_PROJETO','DTESTFECHAMENTO','QTDNEG']].sort_values(by='QTDNEG', ascending=False)

    df = df.merge(tabela_compras, on = 'GRUPO_PRODUTO', how='left')
    df = df.fillna(0)
    df['VLRTOT'] = df['QTDNEG'] * df['VLRUNITDOLAR']

    total_projetos = str(df.shape[0])
    if group == 1:
        df = df.groupby(['GRUPO_PRODUTO']).agg(
            {'NOME_PROJETO': lambda x : '/ '.join(set(str(i) for i in x)),
                'QTDNEG': 'sum','VLRTOT':'sum', 'DTESTFECHAMENTO': 'max'
                }).reset_index().sort_values(by='QTDNEG', ascending=False)

    df = df[['GRUPO_PRODUTO','NOME_PROJETO','QTDNEG','VLRTOT','DTESTFECHAMENTO']]
    df.rename(columns={'GRUPO_ECONOMICO': 'CUSTOMER', 'NOME_PROJETO': 'PROJECT', 'GRUPO_PRODUTO': 'PRODUCT', 'QTDNEG': 'ANUAL VOLUME (KG)', 'VLRTOT': 'VALUE (USD)', 'DTESTFECHAMENTO': 'CONVERSION ESTIMATED DATE'}, inplace=True)
    df['CONVERSION ESTIMATED DATE'] = df['CONVERSION ESTIMATED DATE'].dt.strftime('%d/%m/%Y')
    peso_total = df['ANUAL VOLUME (KG)'].sum()
    dolar_total = df['VALUE (USD)'].sum()
    df['VALUE (USD)'] = df['VALUE (USD)'].apply(lambda x: '${:,.0f}'.format(x)).str.replace(',', '.')
    df['ANUAL VOLUME (KG)'] = df['ANUAL VOLUME (KG)'].apply(lambda x: '{:,.0f}'.format(x)).str.replace(',', '.')
    peso_total = '{:,.0f}'.format(peso_total).replace(',', '.')
    dolar_total = '${:,.0f}'.format(dolar_total).replace(',', '.')

    #-------------------SLIDE-------------------#
    # Get the layout of the slide
    layout_name = 'tabela_simples'

    # Find the slide layout by name and add a new slide
    slide_layout = find_layout_by_name(prs, layout_name)
    slide = prs.slides.add_slide(slide_layout)

    # Set the title for the slide
    text_placeholder = slide.placeholders[0]  # Index 0 for the TITLE placeholder (available in "template_placeholder_mapping.txt)")
    text_placeholder.text = 'Oportunidades Produtos Especificos: ' + 'CARBOPOL EZ 4, NOVERITE 311 POLYMER, NOVERITE LD920N, DILUTHIX CLEAR SOFTENER'

    # Set the summary for the data
    text_placeholder = slide.placeholders[14]  # Index 0 for the TITLE placeholder (available in "template_placeholder_mapping.txt)")
    text_placeholder.text = total_projetos + 'projetos' + '\n' + str(peso_total) + ' Kg' + '\n' + str(dolar_total) + ' USD'

    nova_tabela(df, slide, 23.84,1.5,3.2)

#------------------EXCLUSIVO LUBRIZOL--------------
def oport_ativos (end_date, prs, bu, lim_table, argire): #Oportunidades abertas para a representada "LIPOTEC" no mês vigente (Exceto oportunidades canceladas)
    #----------------CRIAÇÃO DO DATAFRAME DE REFERÊNCIA----------------#
    #Oportunidades com previsão para o ano de referência

    if argire == 1:
        df = df_oportunidade[(df_oportunidade['DESCRICAO_APLICACAO'].str.contains(bu[0])|df_oportunidade['DESCRICAO_APLICACAO'].str.contains(bu[1])) 
                            & (df_oportunidade['GRUPO_PRODUTO'] == 'ARGIRELINE AMPLIFIED PEPTIDE')
                            & (df_oportunidade['STATUS RESUMIDO'] == 'Em Andamento')
                             ][['GRUPO_ECONOMICO','GRUPO_PRODUTO','NOME_PROJETO','DTESTFECHAMENTO','QTDNEG','VLRTOT','STATUS']].sort_values(by='QTDNEG', ascending=False)

    else:
        df = df_oportunidade[df_oportunidade['DESCRICAO_APLICACAO'].str.contains(bu[0])|df_oportunidade['DESCRICAO_APLICACAO'].str.contains(bu[1])]
        df = df[(df['DTNEG'].dt.month == end_date.month) 
                & (df['DTNEG'].dt.year == end_date.year)
                & (df['STATUS RESUMIDO'] != 'Cancelado') 
                & (df['REPRESENTADA'].isin(['LIPOTEC','LIPOTEC USA']))
                ][['GRUPO_ECONOMICO','GRUPO_PRODUTO','NOME_PROJETO','DTESTFECHAMENTO','QTDNEG','VLRTOT','STATUS']].sort_values(by='QTDNEG', ascending=False)

    total_projects = str(df.shape[0])

    df = df.merge(tabela_compras, on = 'GRUPO_PRODUTO', how='left')
    df = df.fillna(0)
    df['VLRTOT'] = df['QTDNEG'] * df['VLRUNITDOLAR']
    df = df[['GRUPO_ECONOMICO','NOME_PROJETO','GRUPO_PRODUTO','STATUS','QTDNEG','VLRTOT','DTESTFECHAMENTO']]
    df.rename(columns={'VLRTOT':'VALUE (USD)','GRUPO_ECONOMICO': 'CUSTOMER', 'NOME_PROJETO': 'PROJECT', 'GRUPO_PRODUTO': 'PRODUCT', 'QTDNEG': 'ANUAL VOLUME (KG)', 'DTESTFECHAMENTO': 'CONVERSION ESTIMATED DATE'}, inplace=True)
    df['CONVERSION ESTIMATED DATE'] = df['CONVERSION ESTIMATED DATE'].dt.strftime('%d/%m/%Y')
    peso_total = df['ANUAL VOLUME (KG)'].sum() 
    dolar_total = df['VALUE (USD)'].sum()
    total_projects = len(df)
    if df.shape[0] > lim_table:
        df = df.head(lim_table)
    if df.shape[0] > 0:
        df.loc['Total', 'ANUAL VOLUME (KG)'] = peso_total
        df.loc['Total', 'VALUE (USD)'] = dolar_total
        df.loc['Total', 'PROJECT'] = f'{str(total_projects)} Projects'
        df.loc['Total', 'CUSTOMER'] = 'Total'
        df['ANUAL VOLUME (KG)'] = df['ANUAL VOLUME (KG)'].apply(lambda x: '{:,.0f}'.format(x)).str.replace(',', '.')
        df['VALUE (USD)'] = df['VALUE (USD)'].apply(lambda x: '$' + '{:,.0f}'.format(x)).str.replace(',', '.')

    #---------------CONFIGURAÇÃO DO SLIDE----------------#
    chart_layout_name = 'tabela_simples_sem_coment'

    # Find the slide layout by name and add a new slide
    slide_layout = find_layout_by_name(prs, chart_layout_name)
    slide = prs.slides.add_slide(slide_layout)

    # Set the title for the slide
    slide.shapes.title.text = ('PROJECTS - OPENED PROJECTS - ACTIVES' if argire == 0 else 'PROJECTS - ARGIRELINE AMPLIFIED')

    # Insert text_summary into the specified text placeholder
    text_placeholder = slide.placeholders[15]  # Index 11 for the text_summary placeholder
    text_placeholder.text = ('New projects (except canceled)' + ' - ' + bu[0] + ' / ' + bu[1] if argire == 0 else 'On going projects' + ' - ' + bu[0] + '/' + bu[1])

    #---------------TABELA----------------#
    individual_width = [5.08,9.5,5.66,3.2,2.71,2.75,3.83]
    nova_tabela(df, slide, individual_width,0.57,3.01) 