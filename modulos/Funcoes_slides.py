from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import pandas as pd

def get_chart_type(chart_type_str):
    chart_types = {
        'bar-horizontal': XL_CHART_TYPE.BAR_CLUSTERED,
        'bar-vertical': XL_CHART_TYPE.COLUMN_CLUSTERED,
        'line': XL_CHART_TYPE.LINE_MARKERS,
        'pie': XL_CHART_TYPE.PIE,
        # Add any other chart types you need
    }
    return chart_types.get(chart_type_str.lower(), XL_CHART_TYPE.COLUMN_CLUSTERED)

# Função para retornar o layout do slide pelo nome
def find_layout_by_name(prs, layout_name):
    for layout in prs.slide_layouts:
        if layout.name == layout_name:
            return layout
    return None

def encurtar_nome(nome):
    df_clientes_tratamento = pd.read_excel('modulos/Parceiro_131124.xlsx', sheet_name='Planilha4')
    df_clientes_tratamento['REMOVER'] = df_clientes_tratamento['REMOVER'].astype(str)
    caracteres_remover = df_clientes_tratamento['REMOVER'].tolist()
    for caracter in caracteres_remover:
        nome = nome.replace(caracter, '')
    return nome

def de_para_representas (nome):
    df_depara_tratamento = pd.read_excel('modulos/Representadas.xlsx', sheet_name='new sheet')
    for i in range(len(df_depara_tratamento)):
        nome = nome.replace(df_depara_tratamento['De'][i], df_depara_tratamento['Para'][i])
    return nome

def nova_tabela (df, slide, individual_width,left,top,posicao_cliente=0, value_size=14):
    # Definir as dimensões da tabela
    rows = len(df)+1
    cols = df.shape[1]
    table_width = Cm(sum(individual_width))
    # Adjust the height of the table according to the number of rows
    table_height = rows

    # Calcular a posição inicial da tabela para centralizá-la no slide
    left = Cm(left)
    top = Cm(top)

    # Adicionar a tabela ao slide
    shape = slide.shapes.add_table(rows,cols,left,top,table_width,table_height)
    table = shape.table

    tbl = shape._element.graphic.graphicData.tbl
    style_id = '{5A111915-BE36-4E01-A7E5-04B1672EAD32}'
    tbl[0][-1].text = style_id

    # Column Heading Format
    for i in range(len(table.columns)):
        table.columns[i].width = Cm(individual_width[i])
        cell = table.cell(0, i)
        cell.text = str(df.columns[i])
        if df.columns[i] == 'Representada':
            representada = i
        for paragraph in cell.text_frame.paragraphs: #Tive que iterar aqui pq tem coluna com "\n", separando em diferentes parágrafos.
            paragraph.font.size = Pt(value_size)
            paragraph.font.name = "Futura Lt BT"
            paragraph.font.color.rgb = RGBColor(250, 250, 250)
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 166, 189)
        # Insert values into the table
        for l in range(rows-1):
            for j in range(cols):
                cell = table.cell(l+1, j)
                cell.text = str(df.iloc[l, j]).upper()
                if j == posicao_cliente:
                    cell.text = encurtar_nome(cell.text)
                try:
                    if j == representada:
                        cell.text = de_para_representas(cell.text)
                except:
                    pass
                pg = cell.text_frame.paragraphs[0]
                # Calculate the width of the text and adjust column width accordingly
                cell.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                pg.font.size = Pt(value_size)
                pg.font.name = "Futura Lt BT"
                pg.alignment = PP_ALIGN.CENTER

def tratamento_grafico_colunas (chart):
    for series in chart.series:
        fill = series.format.fill
        fill.solid()

        # Ajustar largura do espaçamento
        series.format.line.width = Pt(0)  # Remove the border line
        series.gap_width = 100
        # Format the axes
        chart.category_axis.tick_labels.font.size = Pt(14)
        chart.value_axis.tick_labels.font.size = Pt(14)
        # Add data labels formatted as percentages with no decimal points
        series.has_data_labels = True
        for point in series.points:
            point.data_label.number_format = '"#.##0_ ;-#.##0"'
            point.data_label.font.size = Pt(14)
            point.data_label.font.typeface = 'Futura Lt BT'
            point.data_label.font.color.rgb = RGBColor(0x00, 0x00, 0x00)  # Black font color
            
        chart.value_axis.visible = False  # Remove the value axis
        chart.category_axis.visible = True

        # Remove the chart title and legend
        chart.has_title = True
        chart.has_legend = False

    # Remove gridlines for all charts
    if chart.category_axis and chart.category_axis.has_major_gridlines:
        chart.category_axis.major_gridlines.format.line.fill.background()

    if chart.value_axis and chart.value_axis.has_major_gridlines:
        chart.value_axis.major_gridlines.format.line.fill.background()

    # Set the minimum value of the value axis to 0
    chart.value_axis.minimum_scale = 0